"""Geometry + text-fit audit. Flags anything a presenter would have to hand-fix."""
from pptx import Presentation
from pptx.util import Emu

P = r"c:\Data_pulse_kg\DataPulse_Leadership_Deck.pptx"
prs = Presentation(P)
SW = prs.slide_width / 914400
SH = prs.slide_height / 914400
print(f"canvas {SW:.3f} x {SH:.3f}\n")

# Rough advance-width table for Segoe UI at 1pt (em fractions), good enough to
# catch real overflow. Consolas is monospace at 0.55 em.
NARROW = set("iljtfIr.,;:'\"|!()[]{}-` ")
WIDE = set("mwMW@")


def est_width_pt(txt, size, mono=False):
    if mono:
        return len(txt) * size * 0.550
    w = 0.0
    for ch in txt:
        if ch in NARROW:
            w += 0.30
        elif ch in WIDE:
            w += 0.92
        elif ch.isupper() or ch.isdigit():
            w += 0.62
        else:
            w += 0.512
    return w * size


issues = []
for idx, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        L = sh.left / 914400
        T = sh.top / 914400
        W = sh.width / 914400
        H = sh.height / 914400
        R, B = L + W, T + H
        # 1) out of bounds
        if L < -0.01 or T < -0.01 or R > SW + 0.01 or B > SH + 0.01:
            issues.append(f"S{idx:02d} OUT-OF-BOUNDS  {sh.shape_type}  "
                          f"L{L:.2f} T{T:.2f} R{R:.2f} B{B:.2f}  "
                          f"{(sh.text_frame.text[:40] if sh.has_text_frame else '')!r}")
        # 1b) footer intrusion — a content shape running into the footer band.
        # Skips full-bleed backgrounds (H >= 6) and hairline rules (H <= 0.6).
        if 0.6 < H < 6.0 and T < 6.94 < B:
            issues.append(
                f"S{idx:02d} FOOTER-INTRUSION  bottom={B:.2f}in (footer sits at 6.99)  "
                f"{(sh.text_frame.text[:40] if sh.has_text_frame else str(sh.shape_type))!r}")
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        if not tf.text.strip():
            continue
        ml = (tf.margin_left or 0) / 914400
        mr = (tf.margin_right or 0) / 914400
        avail_w_pt = (W - ml - mr) * 72
        # 2) per-paragraph line count vs box height
        total_lines = 0
        max_size = 0
        for p in tf.paragraphs:
            ptxt = "".join(r.text for r in p.runs)
            if not ptxt:
                total_lines += 1
                continue
            sizes = [r.font.size.pt for r in p.runs if r.font.size]
            sz = max(sizes) if sizes else 12
            max_size = max(max_size, sz)
            mono = any((r.font.name or "") == "Consolas" for r in p.runs)
            wpt = est_width_pt(ptxt, sz, mono)
            lines = max(1, int(wpt / avail_w_pt) + (1 if wpt % avail_w_pt else 0))
            # explicit newlines
            lines += ptxt.count("\n")
            total_lines += lines
            # 3) single-line label that cannot fit
            if len(tf.paragraphs) == 1 and H * 72 < sz * 1.45 and wpt > avail_w_pt:
                issues.append(f"S{idx:02d} LABEL-CLIP  h={H:.2f}in size={sz}pt "
                              f"needs {wpt/72:.2f}in has {avail_w_pt/72:.2f}in  {ptxt[:52]!r}")
        ls = tf.paragraphs[0].line_spacing or 1.2
        if isinstance(ls, float):
            need_in = total_lines * max_size * ls * 1.02 / 72
        else:
            need_in = total_lines * (ls.pt if hasattr(ls, "pt") else 14) / 72
        if need_in > H + 0.075:
            issues.append(f"S{idx:02d} TEXT-OVERFLOW  box h={H:.2f}in needs~{need_in:.2f}in "
                          f"({total_lines} lines @ {max_size}pt)  {tf.text[:52]!r}")

print(f"{len(issues)} issue(s)\n")
for i in issues:
    print(i)
