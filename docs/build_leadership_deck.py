"""Build the DataPulse leadership deck — light/white edition.

Clean white canvas, Accenture purple accent retained. Two supporting hues keep
"the graph plans" and "SQLite executes" visually distinct. Every coordinate is
derived from one grid, so the deck is presentation-ready with no manual nudging.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = r"c:\Data_pulse_kg\DataPulse_Leadership_Deck.pptx"

# ── palette — light theme, brand purple preserved ────────────────────────────
BG        = RGBColor(0xFF, 0xFF, 0xFF)   # pure white canvas
PAPER     = RGBColor(0xFA, 0xFA, 0xFC)   # title / closing wash
SURFACE   = RGBColor(0xF6, 0xF6, 0xF9)   # card fill
SURFACE2  = RGBColor(0xEF, 0xEF, 0xF4)   # raised card
BORDER    = RGBColor(0xE1, 0xE1, 0xE9)   # hairline
BORDER_HI = RGBColor(0xC6, 0xC6, 0xD2)
ACCENT    = RGBColor(0xA1, 0x00, 0xFF)   # Accenture purple — bars, ticks, rules
ACCENT_TX = RGBColor(0x6D, 0x00, 0xB3)   # purple that passes contrast as text
KG        = RGBColor(0x6D, 0x00, 0xB3)   # knowledge graph = purple
SQL       = RGBColor(0x0A, 0x6B, 0xAF)   # SQLite = blue
TEXT      = RGBColor(0x11, 0x11, 0x19)
DIM       = RGBColor(0x45, 0x45, 0x4F)
MUTED     = RGBColor(0x78, 0x78, 0x8A)
TEAL      = RGBColor(0x0B, 0x7C, 0x63)
WARN      = RGBColor(0xA6, 0x5E, 0x0A)
ERR       = RGBColor(0xBE, 0x2F, 0x2A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SANS  = "Segoe UI"
LIGHT = "Segoe UI Light"
MONO  = "Consolas"

# ── grid ────────────────────────────────────────────────────────────────────
SW, SH   = 13.333, 7.5
M        = 0.60
CW       = SW - 2 * M
BODY_TOP = 1.78
FOOT_Y   = 6.99

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

_n = [0]


# ── primitives ──────────────────────────────────────────────────────────────

def new_slide(wash: bool = False):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    _style(r, fill=PAPER if wash else BG, line=None)
    return s


def _style(shape, fill=None, line=None, lw=1.0):
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    if shape.has_text_frame:
        shape.text_frame.word_wrap = True
    return shape


def rect(s, x, y, w, h, fill=SURFACE, line=BORDER, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    return _style(s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h)),
                  fill, line, lw)


def tracking(run, pts: float):
    run.font._rPr.set("spc", str(int(pts * 100)))


def text(s, x, y, w, h, spans, size=10.5, color=DIM, bold=False, font=SANS,
         align=PP_ALIGN.LEFT, ls=1.26, anchor=MSO_ANCHOR.TOP, space=0.0, spc=0.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = [spans] if isinstance(spans, str) else spans
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        if space:
            p.space_after = Pt(space)
        runs = [(para, {})] if isinstance(para, str) else para
        runs = [(r, {}) if isinstance(r, str) else r for r in runs]
        for txt, ov in runs:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.color.rgb = ov.get("color", color)
            t = ov.get("spc", spc)
            if t:
                tracking(r, t)
    return tb


def eyebrow(s, x, y, label, color=ACCENT_TX, w=None, size=8.5):
    text(s, x, y, w or 6.0, 0.20, label.upper(), size=size, color=color, bold=True, spc=1.3)


def header(s, kicker, title, sub=None):
    rect(s, M, 0.44, 0.055, 0.30, fill=ACCENT, line=None)
    eyebrow(s, M + 0.17, 0.47, kicker)
    text(s, M, 0.72, CW, 0.48, title, size=25, color=TEXT, bold=True, ls=1.0)
    if sub:
        text(s, M, 1.22, CW - 0.4, 0.26, sub, size=11.5, color=DIM, ls=1.0)
    y = 1.60 if sub else 1.42
    rect(s, M, y, CW, 0.014, fill=ACCENT, line=None)
    return y


def footer(s, label="DataPulse  ·  Knowledge-Graph Grounded Text-to-SQL"):
    _n[0] += 1
    text(s, M, FOOT_Y, 8.0, 0.20, label, size=8, color=MUTED, spc=0.4)
    text(s, SW - M - 2.0, FOOT_Y, 2.0, 0.20, f"{_n[0]:02d}", size=8, color=ACCENT_TX,
         align=PP_ALIGN.RIGHT, bold=True)


def cols(n, gap=0.24, x0=M, total=CW):
    w = (total - gap * (n - 1)) / n
    return [(x0 + i * (w + gap), w) for i in range(n)]


def arrow(s, x, y, w, h, shape=MSO_SHAPE.RIGHT_ARROW, color=BORDER_HI):
    return rect(s, x, y, w, h, fill=color, line=None, shape=shape)


def chip(s, x, y, w, h, label, color=ACCENT_TX, fill=WHITE, size=9.5, mono=False,
         bold=True, line=None):
    r = rect(s, x, y, w, h, fill=fill, line=line or color)
    tf = r.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = label
    run.font.name = MONO if mono else SANS
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return r


def card(s, x, y, w, h, title=None, body=None, accent=ACCENT, fill=SURFACE,
         line=BORDER, tick=True, tsize=12.5, bsize=10.2, tcolor=TEXT, num=None):
    rect(s, x, y, w, h, fill=fill, line=line)
    if tick:
        rect(s, x, y, 0.05, h, fill=accent, line=None)
    px, pw = x + 0.28, w - 0.52
    cy = y + 0.20
    if num:
        text(s, px, cy, pw, 0.24, num, size=9, color=accent, bold=True, spc=1.1, ls=1.0)
        cy += 0.28
    if title:
        text(s, px, cy, pw, 0.30, title, size=tsize, color=tcolor, bold=True, ls=1.05)
        cy += 0.34 if tsize <= 13 else 0.40
    if body:
        text(s, px, cy, pw, y + h - cy - 0.16, body, size=bsize, color=DIM, ls=1.30, space=4)


def kpi(s, x, y, w, h, value, label, note=None, color=ACCENT_TX):
    rect(s, x, y, w, h, fill=SURFACE, line=BORDER)
    rect(s, x, y, w, 0.05, fill=color, line=None)
    text(s, x + 0.24, y + 0.30, w - 0.48, 0.56, value, size=29, color=TEXT, bold=True, ls=1.0)
    text(s, x + 0.24, y + 0.92, w - 0.48, 0.22, label.upper(), size=8.5, color=color,
         bold=True, spc=1.1, ls=1.0)
    if note:
        text(s, x + 0.24, y + 1.16, w - 0.48, 0.44, note, size=9.2, color=MUTED, ls=1.22)


# ═══════════════════════════════════════════════════════════════════════════
# 01 — Title
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(wash=True)
rect(s, 0, 0, SW, 0.14, fill=ACCENT, line=None)
# schema-graph motif, top right
for (nx, ny, r_) in [(11.05, 1.35, 0.16), (12.20, 1.95, 0.11), (10.55, 2.45, 0.11),
                     (11.75, 3.05, 0.14), (10.15, 3.70, 0.10), (12.45, 3.95, 0.09)]:
    rect(s, nx, ny, r_, r_, fill=ACCENT, line=None, shape=MSO_SHAPE.OVAL)
for (a, b, c, d) in [(11.13, 1.51, 12.25, 1.95), (11.13, 1.51, 10.61, 2.45),
                     (12.25, 2.06, 11.82, 3.05), (10.61, 2.56, 11.75, 3.12),
                     (11.75, 3.19, 10.20, 3.70), (11.89, 3.19, 12.49, 3.95)]:
    ln = s.shapes.add_connector(1, Inches(a), Inches(b), Inches(c), Inches(d))
    ln.line.color.rgb = RGBColor(0xD9, 0xB8, 0xF2)
    ln.line.width = Pt(1.25)

rect(s, M, 1.62, 0.055, 0.34, fill=ACCENT, line=None)
eyebrow(s, M + 0.17, 1.66, "Enterprise Data Intelligence  ·  Technical Briefing", w=8.0)

text(s, M, 2.20, 9.6, 0.80, [[("Data", {"font": LIGHT, "color": TEXT}),
                             ("Pulse", {"bold": True, "color": ACCENT_TX})]],
     size=58, ls=0.95)
text(s, M, 3.10, 9.9, 1.10,
     [[("Ask your enterprise data in plain English.", {"color": TEXT, "bold": True})],
      [("A knowledge graph plans the query. The database answers it.", {"color": DIM})]],
     size=19, ls=1.30)

rect(s, M, 4.42, 4.2, 0.014, fill=ACCENT, line=None)
text(s, M, 4.66, 8.4, 0.72,
     "A reasoning agent grounded in a Neo4j schema graph — turning natural-language "
     "business questions into verified SQL, with every step of its reasoning open to inspection.",
     size=12.5, color=DIM, ls=1.42)

for i, (x, w) in enumerate(cols(4, gap=0.20, total=9.6)):
    v, l = [("5", "Domains"), ("50", "Tables"), ("23,293", "Rows"), ("441", "KG Nodes")][i]
    rect(s, x, 5.66, w, 0.86, fill=SURFACE, line=BORDER)
    text(s, x, 5.84, w, 0.34, v, size=19, color=TEXT, bold=True, align=PP_ALIGN.CENTER, ls=1.0)
    text(s, x, 6.20, w, 0.20, l.upper(), size=8, color=ACCENT_TX, bold=True,
         align=PP_ALIGN.CENTER, spc=1.0, ls=1.0)

text(s, M, 6.92, 9.6, 0.24,
     "Neo4j AuraDB   ·   LLaMA-3.3-70B on Groq   ·   FastAPI   ·   SQLite   ·   184 automated tests",
     size=9.5, color=MUTED, spc=0.3)
_n[0] += 1

# ═══════════════════════════════════════════════════════════════════════════
# 02 — Executive summary
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Executive Summary", "What we built, and why it is different",
       "The knowledge graph is not decoration — it is the mechanism that makes LLM-generated SQL trustworthy.")

for i, (x, w) in enumerate(cols(4)):
    v, l, nt, c = [
        ("83.8%", "Prompt context removed", "The graph sends the model 8 relevant tables, not all 50.", ACCENT_TX),
        ("6.2×", "Smaller schema payload", "10,224 → 1,655 tokens per question. Faster, cheaper, sharper.", SQL),
        ("100%", "Read-only enforcement", "Three independent guards. The agent can never write.", TEAL),
        ("184", "Automated tests", "The reasoning loop is unit-tested without a live model.", WARN),
    ][i]
    kpi(s, x, BODY_TOP, w, 1.78, v, l, nt, c)

y2 = BODY_TOP + 2.02
for i, (x, w) in enumerate(cols(3)):
    t, b, n = [
        ("The problem is not SQL generation",
         "Models write plausible SQL easily. They fail at knowing which of 50 tables to use, "
         "how those tables join, and which column the business calls “revenue”. That is a "
         "metadata problem, not a language problem.", "01  DIAGNOSIS"),
        ("The graph supplies the missing map",
         "Neo4j holds the schema — every table, column, foreign key and canonical metric — "
         "searchable by meaning. The agent asks the graph what is relevant, then writes SQL "
         "against a precise, verified slice.", "02  MECHANISM"),
        ("Every answer carries its evidence",
         "The response returns the SQL, the rows, an inspectable reasoning trace, and a "
         "deterministic grounding check that confirms the stated figures actually appear in "
         "the data.", "03  OUTCOME"),
    ][i]
    card(s, x, y2, w, 2.38, t, b, num=n, tsize=13.5)

footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 03 — The problem, told as a story
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "The Problem", "The data is not the bottleneck. The translation is.",
       "Every enterprise already stores the answer. Reaching it requires a person who holds the schema in their head.")

# ── the journey strip: what happens today ──
eyebrow(s, M, BODY_TOP, "The journey of one business question, today", color=MUTED, w=8.0)
jy, jh = BODY_TOP + 0.24, 1.16
stages = [
    ("MON 09:00", "The question", "“What was revenue by region last quarter?”", MUTED),
    ("MON 11:30", "The ticket", "Filed to the analytics queue. Position: seventh.", MUTED),
    ("WED 14:00", "The translation", "An analyst reverse-engineers the schema and the joins.", WARN),
    ("THU 16:00", "The spreadsheet", "A number arrives — with a caveat about which column was used.", WARN),
    ("FRI", "The decision", "Already made, without it.", ERR),
]
jcols = cols(5, gap=0.16)
for i, (when, t, b, c) in enumerate(stages):
    x, w = jcols[i]
    last = i == len(stages) - 1
    rect(s, x, jy, w, jh, fill=SURFACE2 if last else SURFACE,
         line=ERR if last else BORDER, lw=1.4 if last else 1.0)
    rect(s, x, jy, w, 0.045, fill=c, line=None)
    text(s, x + 0.20, jy + 0.18, w - 0.40, 0.20, when, size=7.8, color=c, bold=True, spc=1.0, ls=1.0)
    text(s, x + 0.20, jy + 0.42, w - 0.40, 0.24, t, size=11.6, color=TEXT, bold=True, ls=1.0)
    text(s, x + 0.20, jy + 0.70, w - 0.40, 0.40, b, size=9.0, color=DIM, ls=1.20)
    if not last:
        arrow(s, x + w + 0.015, jy + jh / 2 - 0.065, 0.13, 0.13, color=BORDER_HI)

# ── the four structural causes ──
cy = jy + jh + 0.34
eyebrow(s, M, cy, "Why it keeps happening — four structural causes", color=MUTED, w=8.0)
causes = [
    ("Schema knowledge is tribal", "Which table is authoritative for revenue? The analysts know. "
     "It is written down nowhere the business can read.", KG),
    ("Dashboards answer yesterday", "Pre-built BI covers the questions someone already thought to "
     "ask. Anything genuinely new rejoins the queue.", KG),
    ("The cost is in the joining", "The question is one sentence. The answer is three joins, a date "
     "boundary, and a contested definition.", SQL),
    ("Scale compounds it", "Fifty tables across five domains is a modest footprint. Real estates "
     "run to thousands — and the queue grows with them.", SQL),
]
ccy = cy + 0.24
for i, (t, b, c) in enumerate(causes):
    x, w = cols(4)[i]
    rect(s, x, ccy, w, 1.44, fill=SURFACE, line=BORDER)
    rect(s, x, ccy, 0.05, 1.44, fill=c, line=None)
    text(s, x + 0.26, ccy + 0.20, w - 0.50, 0.44, t, size=11.8, color=TEXT, bold=True, ls=1.06)
    text(s, x + 0.26, ccy + 0.70, w - 0.50, 0.60, b, size=9.4, color=DIM, ls=1.24)

# ── the reframe ──
ry2 = ccy + 1.66
rect(s, M, ry2, CW, 0.66, fill=SURFACE2, line=ACCENT, lw=1.4)
rect(s, M, ry2, 0.055, 0.66, fill=ACCENT, line=None)
text(s, M + 0.30, ry2 + 0.19, CW - 0.70, 0.30,
     [[("The reframe:  ", {"color": ACCENT_TX, "bold": True, "size": 12.5}),
       ("the scarce resource was never the data, or even the SQL. It is the structural knowledge of "
        "how the data fits together — so we made that knowledge a queryable asset.",
        {"color": TEXT, "size": 12.5})]], ls=1.0)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 04 — Why naive LLM text-to-SQL fails
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Why The Obvious Approach Fails", "“Just paste the schema into the prompt” breaks in four predictable ways",
       "Each failure is silent: the model returns confident, well-formed SQL that is wrong.")

fails = [
    ("Context dilution",
     "The whole catalogue is 373 columns and ~10,200 tokens. Buried in it, the model's "
     "attention spreads thin and it picks a plausible-looking wrong column.",
     "50 tables in · 8 relevant"),
    ("Invented joins",
     "Foreign keys are not in the DDL the model sees. It guesses a join key that looks "
     "right, and silently produces a wrong row count.",
     "40 real FK edges · 0 guessable"),
    ("Ambiguous business terms",
     "“Revenue” matches invoices.amount and order_items.line_total. Both are numeric. "
     "Only one is correct — and they differ by 3.8×.",
     "$7,782,964 vs $2,030,281"),
    ("Unverifiable output",
     "A single opaque generation. No trace of why those tables, no check that the stated "
     "number appears anywhere in the result.",
     "Confident · Unauditable"),
]
for i, (t, b, tag) in enumerate(fails):
    x, w = cols(4)[i]
    rect(s, x, BODY_TOP, w, 2.62, fill=SURFACE, line=BORDER)
    rect(s, x, BODY_TOP, w, 0.05, fill=ERR, line=None)
    text(s, x + 0.26, BODY_TOP + 0.26, w - 0.50, 0.24, f"FAILURE {i+1:02d}", size=8.2,
         color=ERR, bold=True, spc=1.1, ls=1.0)
    text(s, x + 0.26, BODY_TOP + 0.56, w - 0.50, 0.30, t, size=13, color=TEXT, bold=True, ls=1.05)
    text(s, x + 0.26, BODY_TOP + 0.98, w - 0.50, 1.10, b, size=10, color=DIM, ls=1.30)
    rect(s, x + 0.26, BODY_TOP + 2.10, w - 0.50, 0.006, fill=BORDER, line=None)
    text(s, x + 0.26, BODY_TOP + 2.24, w - 0.50, 0.24, tag, size=9.4, color=ERR, bold=True,
         font=MONO, ls=1.0)

yb = BODY_TOP + 2.86
rect(s, M, yb, CW, 1.58, fill=SURFACE2, line=ACCENT, lw=1.4)
rect(s, M, yb, 0.055, 1.58, fill=ACCENT, line=None)
text(s, M + 0.30, yb + 0.24, CW - 0.70, 0.26, "THE COMMON ROOT CAUSE", size=8.6,
     color=ACCENT_TX, bold=True, spc=1.2, ls=1.0)
text(s, M + 0.30, yb + 0.58, CW - 0.70, 0.80,
     [[("All four failures are the same failure: the model is asked to plan a query without a map of the schema. ",
        {"color": TEXT, "bold": True, "size": 14}),
       ("Fixing it does not need a bigger model or a longer prompt — it needs the structural knowledge "
        "(what joins to what, what a metric means, what is relevant to this question) to be stored, "
        "searchable, and handed to the model as fact rather than left to inference.",
        {"color": DIM, "size": 12.4})]],
     ls=1.36)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 05 — Two databases, two jobs
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "The Core Design Decision", "Two databases, two jobs",
       "The single most important architectural choice in the system — and the reason the graph earns its place.")

bw = (CW - 0.90) / 2
box_h = 3.00
for i, (t, sub, rows, c) in enumerate([
    ("Neo4j — the Knowledge Graph", "SOURCE OF THE PLAN",
     [("Holds", "Schema metadata only — never row data"),
      ("Nodes", "5 Domain · 50 Table · 373 Column · 13 Metric"),
      ("Edges", "40 :REFERENCES join edges carrying exact keys"),
      ("Search", "3 cosine vector indexes over 384-d embeddings"),
      ("Answers", "“Which tables matter, and how do they join?”")], KG),
    ("SQLite — the Execution Store", "SOURCE OF THE ANSWER",
     [("Holds", "The actual business rows — 23,293 across 50 tables"),
      ("Access", "Read-only connection, 5 s timeout, 200-row cap"),
      ("Guard", "SELECT / WITH allowlist; writes and DDL rejected"),
      ("Used by", "run_sql and sample_values tools only"),
      ("Answers", "“What is the number?”")], SQL),
]):
    x = M + i * (bw + 0.90)
    rect(s, x, BODY_TOP, bw, box_h, fill=SURFACE, line=BORDER)
    rect(s, x, BODY_TOP, bw, 0.05, fill=c, line=None)
    text(s, x + 0.30, BODY_TOP + 0.28, bw - 0.60, 0.22, sub, size=8.4, color=c, bold=True,
         spc=1.2, ls=1.0)
    text(s, x + 0.30, BODY_TOP + 0.56, bw - 0.60, 0.32, t, size=15, color=TEXT, bold=True, ls=1.0)
    rect(s, x + 0.30, BODY_TOP + 1.00, bw - 0.60, 0.006, fill=BORDER, line=None)
    for j, (k, v) in enumerate(rows):
        ry = BODY_TOP + 1.16 + j * 0.36
        text(s, x + 0.30, ry, 0.90, 0.22, k.upper(), size=8.2, color=MUTED, bold=True, spc=0.8, ls=1.0)
        text(s, x + 1.26, ry - 0.025, bw - 1.56, 0.30, v, size=9.9, color=DIM, ls=1.16)

cx = M + bw + 0.45
rect(s, cx - 0.007, BODY_TOP + 0.30, 0.014, box_h - 0.60, fill=BORDER_HI, line=None)
chip(s, cx - 0.30, BODY_TOP + box_h / 2 - 0.16, 0.60, 0.32, "vs", color=MUTED, size=9,
     line=BORDER_HI)

yb = BODY_TOP + box_h + 0.30
for i, (x, w) in enumerate(cols(3)):
    t, b = [
        ("Regenerating data never rebuilds the graph",
         "The graph is metadata-only, so changing the generator seed leaves it untouched. Build once, query forever."),
        ("The plan is auditable before a row is read",
         "Retrieval returns the chosen tables and join keys as data — reviewable independently of any SQL."),
        ("Each store is guarded for its own risk",
         "The graph is read-only by nature. SQLite is read-only by enforcement. Neither can be mutated by a query."),
    ][i]
    rect(s, x, yb, w, 1.10, fill=SURFACE2, line=BORDER)
    text(s, x + 0.26, yb + 0.20, w - 0.50, 0.24, t, size=11.2, color=TEXT, bold=True, ls=1.05)
    text(s, x + 0.26, yb + 0.54, w - 0.50, 0.44, b, size=9.5, color=DIM, ls=1.26)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 06 — BLOCK DIAGRAM  (layered system at a glance)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "System Block Diagram", "The whole system in six layers")

LX, LW = M, 1.52                 # layer-label column
BXX = M + LW + 0.20              # blocks column
BXW = CW - LW - 0.20
BAND_H, BAND_GAP = 0.72, 0.145
BY0 = 1.60

bands = [
    ("Presentation", ACCENT, [
        ("Browser UI", "vanilla JS · zero framework"),
        ("Live pipeline + trace", "every step rendered as it happens"),
        ("Data-quality panel", "profile of the generated dataset")]),
    ("API layer", ACCENT, [
        ("POST /api/query", "question → grounded answer"),
        ("GET /api/health", "graph reachability + freshness"),
        ("GET /api/domains", "catalogue domains"),
        ("GET /api/quality", "validator verdict + profile")]),
    ("Reasoning", ACCENT, [
        ("ReAct agent loop", "plan → act → observe · max 6 steps"),
        ("LLaMA-3.3-70B · Groq", "tool-calling inference, temp 0.1"),
        ("Grounding check", "figures matched to rows, deterministic")]),
    ("Tool layer", KG, [
        ("get_schema_context", "→ Neo4j · tables, join keys, metrics"),
        ("sample_values", "→ SQLite · resolve a categorical filter"),
        ("run_sql", "→ SQLite · read-only SELECT / WITH")]),
    ("Knowledge & data", KG, [
        ("Neo4j AuraDB", "441 nodes · 40 join edges · 3 vector indexes"),
        ("all-MiniLM-L6-v2", "384-d embeddings, local, no API cost"),
        ("SQLite  sales.db", "50 tables · 23,293 rows · opened read-only")]),
    ("Catalogue & build", SQL, [
        ("schema.json", "the single contract every layer reads"),
        ("datagen", "seeded CSVs, one module per domain"),
        ("quality gate", "conformance + referential integrity"),
        ("loader + KG builder", "→ SQLite  ·  → Neo4j, idempotent")]),
]

for bi, (layer, lc, blocks) in enumerate(bands):
    by = BY0 + bi * (BAND_H + BAND_GAP)
    # layer label — number and name in one block, vertically centred (no collision)
    rect(s, LX, by, LW, BAND_H, fill=SURFACE2, line=BORDER)
    rect(s, LX, by, 0.05, BAND_H, fill=lc, line=None)
    text(s, LX + 0.18, by, LW - 0.32, BAND_H,
         [[(f"L{6 - bi}", {"color": MUTED, "size": 8.2}),
           ("   " + layer, {"color": lc, "size": 10.4, "bold": True})]],
         ls=1.10, anchor=MSO_ANCHOR.MIDDLE)
    # blocks
    bcols = cols(len(blocks), gap=0.16, x0=BXX, total=BXW)
    for i, (t, sub) in enumerate(blocks):
        x, w = bcols[i]
        rect(s, x, by, w, BAND_H, fill=SURFACE, line=BORDER)
        rect(s, x, by, w, 0.04, fill=lc, line=None)
        mono = any(ch in t for ch in "_/.") and " " not in t.strip()
        text(s, x + 0.18, by + 0.17, w - 0.34, 0.24, t, size=10.4, color=TEXT, bold=True,
             font=MONO if mono else SANS, ls=1.0)
        text(s, x + 0.18, by + 0.43, w - 0.34, 0.26, sub, size=8.4, color=MUTED, ls=1.12)
    # inter-band arrow
    if bi < len(bands) - 1:
        arrow(s, BXX + 0.44, by + BAND_H + 0.012, 0.13, 0.12, MSO_SHAPE.DOWN_ARROW,
              color=BORDER_HI)

text(s, BXX, BY0 + 6 * (BAND_H + BAND_GAP) - 0.02, BXW, 0.22,
     [[("Request flows down L6 → L1;  the grounded response returns up with its SQL, rows and trace.  ",
        {"color": DIM, "size": 9.2}),
       ("One FastAPI process serves every layer above the datastores.",
        {"color": ACCENT_TX, "size": 9.2, "bold": True})]], ls=1.0)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 07 — End-to-end architecture (runtime flow)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "End-to-End Architecture", "One catalogue contract, two stores, one reasoning loop")

BT = 1.62
eyebrow(s, M, BT, "① Build time  ·  offline, idempotent, re-runnable", color=MUTED, w=8.0)
by, bh = BT + 0.24, 0.62
bcols = [(M + i * (2.72 + 0.42), 2.72) for i in range(4)]
build = [("schema.json", "the catalogue contract"),
         ("datagen", "one seeded CSV per table"),
         ("quality gate", "conformance + integrity"),
         ("load", "→ SQLite  ·  → Neo4j")]
for i, (t, sub) in enumerate(build):
    x, w = bcols[i]
    rect(s, x, by, w, bh, fill=SURFACE, line=BORDER)
    rect(s, x, by, 0.05, bh, fill=ACCENT, line=None)
    text(s, x + 0.22, by + 0.11, w - 0.36, 0.22, t, size=11.4, color=TEXT, bold=True,
         font=MONO, ls=1.0)
    text(s, x + 0.22, by + 0.36, w - 0.36, 0.20, sub, size=8.6, color=MUTED, ls=1.0)
    if i < 3:
        arrow(s, x + w + 0.10, by + bh / 2 - 0.075, 0.22, 0.15)

sy = by + bh + 0.30
eyebrow(s, M, sy, "② Shared state  ·  two stores, two jobs", color=MUTED, w=8.0)
stz, sth = sy + 0.24, 0.80
sw_ = (CW - 0.36) / 2
for i, (t, sub, c) in enumerate([
    ("Neo4j  ·  Knowledge Graph", "441 nodes  ·  40 join edges  ·  3 vector indexes (384-d)", KG),
    ("SQLite  ·  sales.db", "50 tables  ·  23,293 rows  ·  opened read-only", SQL),
]):
    x = M + i * (sw_ + 0.36)
    rect(s, x, stz, sw_, sth, fill=SURFACE2, line=c, lw=1.4)
    text(s, x + 0.28, stz + 0.14, sw_ - 0.56, 0.26, t, size=12.8, color=c, bold=True, ls=1.0)
    text(s, x + 0.28, stz + 0.45, sw_ - 0.56, 0.22, sub, size=9.2, color=DIM, ls=1.0)
    chip(s, x + sw_ - 1.30, stz + 0.16, 1.02, 0.24,
         "THE PLAN" if i == 0 else "THE DATA", color=c, size=7.8)
arrow(s, bcols[3][0] + 1.30, by + bh + 0.04, 0.15, 0.20, MSO_SHAPE.DOWN_ARROW)

ty = stz + sth + 0.34
eyebrow(s, M, ty, "③ Query time  ·  the three tools the model may call", color=MUTED, w=9.0)
tz, th = ty + 0.24, 0.66
tools = [("get_schema_context", "tables · exact join keys · metric defs", KG, 4.30),
         ("sample_values", "distinct values → resolve a filter", SQL, 3.72),
         ("run_sql", "execute the SELECT / WITH", SQL, 3.75)]
tx = M
for t, sub, c, w in tools:
    rect(s, tx, tz, w, th, fill=SURFACE, line=BORDER)
    rect(s, tx, tz, 0.05, th, fill=c, line=None)
    text(s, tx + 0.22, tz + 0.12, w - 0.36, 0.22, t, size=10.8, color=c, bold=True,
         font=MONO, ls=1.0)
    text(s, tx + 0.22, tz + 0.38, w - 0.36, 0.20, sub, size=8.6, color=MUTED, ls=1.0)
    tx += w + 0.18
arrow(s, M + 2.05, stz + sth + 0.05, 0.15, 0.24, MSO_SHAPE.UP_DOWN_ARROW, KG)
arrow(s, M + 7.20, stz + sth + 0.05, 0.15, 0.24, MSO_SHAPE.UP_DOWN_ARROW, SQL)

ay = tz + th + 0.32
alh = 0.86
q_w, ag_w = 2.30, 5.10
ans_w = CW - q_w - ag_w - 0.72
rect(s, M, ay, q_w, alh, fill=SURFACE, line=BORDER)
text(s, M + 0.22, ay + 0.16, q_w - 0.44, 0.24, "Question", size=11.6, color=TEXT, bold=True, ls=1.0)
text(s, M + 0.22, ay + 0.45, q_w - 0.44, 0.32, "plain English, from\nthe browser", size=8.8,
     color=MUTED, ls=1.16)
arrow(s, M + q_w + 0.10, ay + alh / 2 - 0.075, 0.20, 0.15, color=ACCENT)

axx = M + q_w + 0.40
rect(s, axx, ay, ag_w, alh, fill=SURFACE2, line=ACCENT, lw=1.5)
rect(s, axx, ay, 0.05, alh, fill=ACCENT, line=None)
text(s, axx + 0.28, ay + 0.13, ag_w - 0.56, 0.26,
     [[("ReAct agent loop", {"color": TEXT, "bold": True, "size": 13.2}),
       ("   ·   LLaMA-3.3-70B on Groq", {"color": ACCENT_TX, "size": 10})]], ls=1.0)
text(s, axx + 0.28, ay + 0.45, ag_w - 0.56, 0.34,
     "Calls tools → reads observations → refines. Up to 6 steps, then a final answer.",
     size=9.2, color=DIM, ls=1.16)

arrow(s, axx + ag_w + 0.10, ay + alh / 2 - 0.075, 0.20, 0.15, color=ACCENT)
rxx = axx + ag_w + 0.40
rect(s, rxx, ay, ans_w, alh, fill=SURFACE, line=TEAL, lw=1.4)
text(s, rxx + 0.22, ay + 0.14, ans_w - 0.44, 0.24, "Grounded response", size=11.6,
     color=TEAL, bold=True, ls=1.0)
text(s, rxx + 0.22, ay + 0.43, ans_w - 0.44, 0.36,
     "answer · sql · rows · trace[] · grounded flag", size=8.8, color=DIM, ls=1.16, font=MONO)
arrow(s, axx + 1.30, tz + th + 0.03, 0.15, 0.26, MSO_SHAPE.UP_DOWN_ARROW, ACCENT)
arrow(s, axx + 3.40, tz + th + 0.03, 0.15, 0.26, MSO_SHAPE.UP_DOWN_ARROW, ACCENT)

fy = ay + alh + 0.24
rect(s, M, fy, CW, 0.40, fill=SURFACE2, line=BORDER)
text(s, M + 0.26, fy + 0.10, CW - 0.52, 0.24,
     [[("Served by one FastAPI process  ", {"color": DIM, "size": 9.4}),
       ("src/api/main.py", {"color": ACCENT_TX, "size": 9.4, "font": MONO}),
       ("  — the API and the UI share a single origin. Deployed on Railway, auto-releasing on every push.",
        {"color": DIM, "size": 9.4})]], ls=1.0)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 08 — TOOLS & TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Tools & Technology", "Every component, and the job it does",
       "Managed services where they earn their cost; local and zero-infrastructure everywhere else.")

groups = [
    ("Knowledge graph", KG, [
        ("Neo4j AuraDB", "Managed graph database — the schema store"),
        ("neo4j Python driver 6.0.2", "Bounded, pooled sessions; health probes"),
        ("Cypher", "shortestPath join discovery · vector queryNodes"),
        ("Native vector indexes", "3 cosine indexes: Column, Table, Domain"),
    ]),
    ("AI & inference", ACCENT_TX, [
        ("Groq API", "Low-latency inference host for tool-calling"),
        ("LLaMA-3.3-70B-versatile", "Plans the query, writes SQL, synthesises"),
        ("groq SDK 1.5.0", "Retry-After handling · malformed-call recovery"),
        ("ReAct tool-calling loop", "Our own loop — provider-agnostic by design"),
    ]),
    ("Embeddings & retrieval", SQL, [
        ("sentence-transformers 4.1.0", "Local encoder — no per-query API cost"),
        ("all-MiniLM-L6-v2", "384-dimension L2-normalised vectors"),
        ("Soft domain routing", "Ranks domains without hard-filtering them"),
        ("Join-path expansion", "Shortest paths over :REFERENCES, ≤ 3 hops"),
    ]),
    ("Backend & API", ACCENT_TX, [
        ("FastAPI 0.115.9", "Typed endpoints, one origin for API and UI"),
        ("Uvicorn 0.34.3", "ASGI server, the production process"),
        ("Pydantic", "Request and response contracts"),
        ("python-dotenv", "Credential and configuration loading"),
    ]),
    ("Data & quality", SQL, [
        ("SQLite 3", "Execution store, opened read-only per query"),
        ("pandas 2.3.0", "CSV generation and bulk load"),
        ("Faker", "Seeded synthetic records — seed 42, deterministic"),
        ("Schema-driven validator", "Conformance + referential integrity gate"),
    ]),
    ("Delivery & assurance", TEAL, [
        ("Vanilla JS + CSS", "Zero-dependency frontend, no build step"),
        ("pytest", "184 tests — the loop runs without a live model"),
        ("Railway", "Health-checked auto-deploy, config as code"),
        ("GitHub", "Source of truth; every push triggers a release"),
    ]),
]

GH = 2.30
for gi, (name, c, rows) in enumerate(groups):
    x, w = cols(3, gap=0.26)[gi % 3]
    y = BODY_TOP + (gi // 3) * (GH + 0.24)
    rect(s, x, y, w, GH, fill=SURFACE, line=BORDER)
    rect(s, x, y, w, 0.05, fill=c, line=None)
    text(s, x + 0.26, y + 0.22, w - 0.50, 0.24, name.upper(), size=9, color=c, bold=True,
         spc=1.15, ls=1.0)
    rect(s, x + 0.26, y + 0.52, w - 0.52, 0.006, fill=BORDER, line=None)
    for ri, (tool, role) in enumerate(rows):
        ry = y + 0.62 + ri * 0.40
        text(s, x + 0.26, ry, w - 0.50, 0.20, tool, size=9.6, color=TEXT, bold=True, ls=1.0)
        text(s, x + 0.26, ry + 0.185, w - 0.50, 0.20, role, size=8.4, color=MUTED, ls=1.0)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 09 — Inside the knowledge graph
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Inside The Knowledge Graph", "What we actually store in Neo4j",
       "Schema metadata, modelled as a graph — because the questions we ask of it are traversal questions.")

lw2 = 6.90
card(s, M, BODY_TOP, lw2, 2.66, "The node model", None, accent=KG)
nodes = [("Domain", "5", "name · description · embedding"),
         ("Table", "50", "name · description · domain · embedding"),
         ("Column", "373", "type · PK flag · aliases · allowed values · embedding"),
         ("Metric", "13", "canonical SQL expression · owning tables")]
for i, (n, cnt, props) in enumerate(nodes):
    ry = BODY_TOP + 0.68 + i * 0.47
    chip(s, M + 0.30, ry, 1.42, 0.34, f":{n}", color=KG, size=9.8, mono=True)
    text(s, M + 2.00, ry + 0.02, 0.66, 0.24, cnt, size=12.4, color=TEXT, bold=True, ls=1.0)
    text(s, M + 2.82, ry + 0.075, lw2 - 3.18, 0.24, props, size=9.3, color=DIM, ls=1.0)

rx2 = M + lw2 + 0.36
rw2 = CW - lw2 - 0.36
card(s, rx2, BODY_TOP, rw2, 2.66, "The relationships", None, accent=ACCENT)
rels = [(":HAS_TABLE", "Domain → Table", "domain membership"),
        (":HAS_COLUMN", "Table → Column", "schema composition"),
        (":FOREIGN_KEY", "Column → Column", "column-level FK"),
        (":REFERENCES", "Table → Table", "carries from_column, to_column")]
for i, (r_, sig, note) in enumerate(rels):
    ry = BODY_TOP + 0.68 + i * 0.47
    hot = r_ == ":REFERENCES"
    text(s, rx2 + 0.28, ry + 0.02, 1.70, 0.24, r_, size=9.8,
         color=ACCENT_TX if hot else DIM, bold=True, font=MONO, ls=1.0)
    text(s, rx2 + 2.06, ry + 0.02, 1.60, 0.24, sig, size=9.5,
         color=TEXT if hot else DIM, bold=hot, ls=1.0)
    text(s, rx2 + 0.28, ry + 0.24, rw2 - 0.56, 0.22, note, size=8.6, color=MUTED, ls=1.0)

yb = BODY_TOP + 2.90
rect(s, M, yb, CW, 1.54, fill=SURFACE2, line=ACCENT, lw=1.4)
rect(s, M, yb, 0.055, 1.54, fill=ACCENT, line=None)
text(s, M + 0.30, yb + 0.22, CW - 0.70, 0.26,
     ":REFERENCES IS THE EDGE THAT DOES THE WORK", size=8.6, color=ACCENT_TX, bold=True,
     spc=1.2, ls=1.0)
text(s, M + 0.30, yb + 0.54, 7.30, 0.84,
     [[("Each of the 40 foreign keys is projected onto a table-to-table edge that carries the "
        "exact join columns as edge properties. ", {"color": DIM, "size": 11}),
       ("That single modelling choice turns “how do these tables join?” — a question an LLM can "
        "only guess at — into a shortest-path query with a provably correct answer.",
        {"color": TEXT, "size": 11, "bold": True})]], ls=1.32)
for i, (x, w) in enumerate(cols(3, gap=0.16, x0=M + 7.90, total=CW - 7.90 - 0.24)):
    v, l = [("3", "vector indexes"), ("384", "dimensions"), ("2", "self-joins")][i]
    rect(s, x, yb + 0.36, w, 0.80, fill=WHITE, line=BORDER)
    text(s, x, yb + 0.50, w, 0.30, v, size=17, color=TEXT, bold=True, align=PP_ALIGN.CENTER, ls=1.0)
    text(s, x, yb + 0.84, w, 0.20, l.upper(), size=7.6, color=ACCENT_TX, bold=True,
         align=PP_ALIGN.CENTER, spc=0.7, ls=1.0)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 10 — Retrieval
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Retrieval  ·  How The Graph Plans", "From a sentence to a precise, joinable schema slice — in five steps",
       "This runs before the model writes a single character of SQL.")

steps = [
    ("01", "Embed the question", "The sentence becomes a 384-dimension vector using a local "
     "all-MiniLM-L6-v2 model. No API call, no cost, runs inside the container."),
    ("02", "Route the domain", "Rank the 5 :Domain nodes by vector similarity. Deliberately a "
     "soft ranking, never a hard filter — cross-domain questions must still work."),
    ("03", "Seed the tables", "Vector-search :Column and :Table nodes, with a recall boost "
     "scoped to the top-ranked domains. Keep the best 8 tables as seeds."),
    ("04", "Expand the join paths", "Walk shortestPath over :REFERENCES between every pair of "
     "seeds, up to 3 hops, and add self-referential loops. Returns the exact join keys."),
    ("05", "Assemble the context", "Fetch columns, types, PKs and allowed values for the final "
     "table set, plus any canonical metric whose tables intersect it."),
]
sh_ = 0.90
for i, (n, t, b) in enumerate(steps):
    yy = BODY_TOP + i * (sh_ + 0.11)
    rect(s, M, yy, 8.10, sh_, fill=SURFACE, line=BORDER)
    rect(s, M, yy, 0.05, sh_, fill=ACCENT, line=None)
    text(s, M + 0.28, yy + 0.26, 0.46, 0.30, n, size=15, color=ACCENT_TX, bold=True,
         font=LIGHT, ls=1.0)
    text(s, M + 0.92, yy + 0.16, 6.88, 0.26, t, size=12.4, color=TEXT, bold=True, ls=1.0)
    text(s, M + 0.92, yy + 0.46, 6.96, 0.40, b, size=9.6, color=DIM, ls=1.24)

ox = M + 8.10 + 0.36
ow = CW - 8.10 - 0.36
card(s, ox, BODY_TOP, ow, 2.42, "What retrieval returns", None, accent=SQL)
outs = [("tables", "descriptions, columns, types, PKs, allowed values"),
        ("joins", "exact from/to columns + cardinality hint"),
        ("domains", "the soft routing ranking, with scores"),
        ("metrics", "canonical expressions for the business terms"),
        ("cross_domain_unjoinable", "a flag that stops invented joins")]
for i, (k, v) in enumerate(outs):
    ry = BODY_TOP + 0.62 + i * 0.36
    text(s, ox + 0.28, ry, ow - 0.56, 0.20, k, size=9.4, color=SQL, bold=True, font=MONO, ls=1.0)
    text(s, ox + 0.28, ry + 0.17, ow - 0.56, 0.20, v, size=8.7, color=MUTED, ls=1.0)

card(s, ox, BODY_TOP + 2.62, ow, 2.42, "Why a graph, not a lookup table",
     [[("A flat metadata table can tell you a foreign key exists. It cannot tell you that "
        "customers reaches order_items through orders in two hops.", {})],
      [("Join discovery is inherently a traversal problem — so we store the schema in the one "
        "database built for traversal.", {"color": ACCENT_TX, "bold": True})]],
     accent=ACCENT, bsize=10.0)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 11 — The reasoning loop
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "The Reasoning Loop", "The model does not answer. It investigates, then answers.",
       "A ReAct tool-calling loop — the reasoning becomes an explicit sequence of steps rather than one opaque generation.")

lwd, lhd = 7.05, 3.66
rect(s, M, BODY_TOP, lwd, lhd, fill=SURFACE, line=BORDER)
text(s, M + 0.30, BODY_TOP + 0.22, lwd - 0.60, 0.26, "THE LOOP  ·  max 6 steps", size=8.6,
     color=ACCENT_TX, bold=True, spc=1.2, ls=1.0)

lx0 = M + 0.42
node_w = lwd - 0.84
seq = [("Send messages to the model", "system prompt + question + every prior observation", ACCENT),
       ("Did the reply contain tool calls?", "the single decision that drives the whole loop", WARN),
       ("Yes → execute each tool, record a Step", "an exception becomes {\"error\": …} and is fed back, never a crash", KG),
       ("No → the reply text is the final answer", "stopped = \"final\"", TEAL)]
for i, (t, b, c) in enumerate(seq):
    yy = BODY_TOP + 0.62 + i * 0.70
    rect(s, lx0, yy, node_w, 0.54, fill=WHITE, line=BORDER)
    rect(s, lx0, yy, 0.05, 0.54, fill=c, line=None)
    text(s, lx0 + 0.24, yy + 0.07, node_w - 0.48, 0.24, t, size=11, color=TEXT, bold=True, ls=1.0)
    text(s, lx0 + 0.24, yy + 0.31, node_w - 0.48, 0.20, b, size=8.7, color=MUTED, ls=1.0)
    if i < 3:
        arrow(s, lx0 + 0.32, yy + 0.555, 0.12, 0.115, MSO_SHAPE.DOWN_ARROW)
text(s, lx0, BODY_TOP + 3.32, node_w, 0.22,
     "Otherwise, loop back to the top — with every observation appended to the conversation.",
     size=8.8, color=ACCENT_TX, bold=True, ls=1.0)

rx3 = M + lwd + 0.36
rw3 = CW - lwd - 0.36
card(s, rx3, BODY_TOP, rw3, 1.72, "One agent, two faces",
     [[("run_agent", {"font": MONO, "color": ACCENT_TX, "bold": True}),
       ("  — the pure loop. Imports no Groq, no Neo4j, no SQLite. Tested with a scripted fake model.", {})],
      [("answer_question", {"font": MONO, "color": ACCENT_TX, "bold": True}),
       ("  — the wiring. Binds the real model, the graph-backed tools, and the catalogue-derived prompt.", {})]],
     accent=ACCENT, bsize=9.8)
card(s, rx3, BODY_TOP + 1.92, rw3, 1.74, "Why that seam matters",
     [[("Changing LLM provider replaces one callable. Nothing in the loop changes — and the loop "
        "is the part that must not regress.", {})],
      [("It is also why 184 tests can cover the reasoning logic with no live model, no credentials "
        "and no flakiness.", {"color": ACCENT_TX, "bold": True})]],
     accent=TEAL, bsize=9.8)

yb = BODY_TOP + 3.80
gr = [("Read-only SQL", "SELECT / WITH allowlist · single statement · writes and DDL rejected"),
      ("Read-only engine", "connection opened ?mode=ro · 5 s timeout · 200-row cap"),
      ("Injection-safe", "table and column names validated against schema.json before quoting"),
      ("Self-healing", "a rejection returns a structured error the agent reads and reforms")]
for i, (t, b) in enumerate(gr):
    x, w = cols(4)[i]
    rect(s, x, yb, w, 1.06, fill=SURFACE2, line=BORDER)
    rect(s, x, yb, w, 0.04, fill=TEAL, line=None)
    text(s, x + 0.24, yb + 0.18, w - 0.48, 0.24, t, size=10.8, color=TEXT, bold=True, ls=1.0)
    text(s, x + 0.24, yb + 0.48, w - 0.48, 0.46, b, size=9.0, color=DIM, ls=1.22)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 12 — Optimisation 1: context economy
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Optimisation 01  ·  Inference", "The graph makes the model's job smaller — measurably",
       "Retrieval replaces the whole catalogue with the slice that matters. Measured on this repository's schema.")

bx, bw2 = M, 7.55
rect(s, bx, BODY_TOP, bw2, 2.96, fill=SURFACE, line=BORDER)
text(s, bx + 0.30, BODY_TOP + 0.24, bw2 - 0.60, 0.26,
     "SCHEMA PAYLOAD SENT TO THE MODEL, PER QUESTION", size=8.6, color=ACCENT_TX, bold=True,
     spc=1.15, ls=1.0)

BAR_X, BAR_MAX = bx + 0.30, bw2 - 2.30
for i, (lab, sub, tokens, frac, c) in enumerate([
    ("Whole catalogue in the prompt", "50 tables · 373 columns · 40,897 characters", "10,224", 1.0, ERR),
    ("Graph-retrieved slice", "8 tables · ~59 columns · ~6,621 characters", "1,655", 1655 / 10224, TEAL),
]):
    yy = BODY_TOP + 0.72 + i * 1.06
    text(s, BAR_X, yy, BAR_MAX, 0.22, lab, size=10.8, color=TEXT, bold=True, ls=1.0)
    text(s, BAR_X, yy + 0.23, BAR_MAX, 0.20, sub, size=8.7, color=MUTED, ls=1.0)
    rect(s, BAR_X, yy + 0.50, BAR_MAX, 0.26, fill=SURFACE2, line=BORDER)
    rect(s, BAR_X, yy + 0.50, max(BAR_MAX * frac, 0.10), 0.26, fill=c, line=None)
    text(s, BAR_X + BAR_MAX + 0.16, yy + 0.44, 1.60, 0.28,
         [[(tokens, {"size": 15, "color": TEXT, "bold": True}),
           ("  tok", {"size": 9, "color": MUTED})]], ls=1.0)

rect(s, bx + 0.30, BODY_TOP + 2.58, bw2 - 0.60, 0.006, fill=BORDER, line=None)
text(s, bx + 0.30, BODY_TOP + 2.68, bw2 - 0.60, 0.22,
     [[("83.8% of the schema context removed  ·  6.2× smaller payload", {"color": TEAL, "bold": True, "size": 10}),
       ("     token estimate = characters ÷ 4", {"color": MUTED, "size": 8.6})]], ls=1.0)

ox2 = M + bw2 + 0.36
ow2 = CW - bw2 - 0.36
gains = [("Higher accuracy", "Fewer candidate columns means fewer ways to be plausibly wrong. Attention lands on relevant schema.", ACCENT),
         ("Lower latency", "A smaller prompt is a faster first token — and the loop may run up to 6 model calls per question.", SQL),
         ("Lower cost", "Token spend falls with the payload, on every call in every step of every question.", TEAL),
         ("Room to grow", "At 500 tables the naive prompt stops fitting. Retrieval returns 8 tables either way.", WARN)]
for i, (t, b, c) in enumerate(gains):
    yy = BODY_TOP + i * 0.76
    rect(s, ox2, yy, ow2, 0.70, fill=SURFACE, line=BORDER)
    rect(s, ox2, yy, 0.05, 0.70, fill=c, line=None)
    text(s, ox2 + 0.24, yy + 0.09, ow2 - 0.48, 0.22, t, size=10.8, color=TEXT, bold=True, ls=1.0)
    text(s, ox2 + 0.24, yy + 0.33, ow2 - 0.48, 0.32, b, size=8.8, color=DIM, ls=1.20)

yb = BODY_TOP + 3.14
rect(s, M, yb, CW, 1.28, fill=SURFACE2, line=ACCENT, lw=1.4)
rect(s, M, yb, 0.055, 1.28, fill=ACCENT, line=None)
text(s, M + 0.30, yb + 0.20, CW - 0.70, 0.26, "THE SCALING ARGUMENT", size=8.6,
     color=ACCENT_TX, bold=True, spc=1.2, ls=1.0)
text(s, M + 0.30, yb + 0.52, CW - 0.70, 0.62,
     [[("The naive approach degrades as the estate grows — more tables means a longer prompt, "
        "more distractors and worse SQL. ", {"color": DIM, "size": 11.4}),
       ("Graph retrieval is flat: it returns the top 8 tables whether the catalogue holds 50 or "
        "5,000. The cost of scale moves out of the prompt and into an indexed vector search.",
        {"color": TEXT, "size": 11.4, "bold": True})]], ls=1.30)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 13 — Optimisation 2: semantic layer
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Optimisation 02  ·  Correctness", "The semantic layer: what the business means by “revenue”",
       "Two columns are numeric, plausible and named like money. Only one is revenue — and the graph knows which.")

tw, thh = 7.55, 2.56
rect(s, M, BODY_TOP, tw, thh, fill=SURFACE, line=BORDER)
text(s, M + 0.30, BODY_TOP + 0.24, tw - 0.60, 0.26, "A LIVE AMBIGUITY IN THIS SCHEMA",
     size=8.6, color=ERR, bold=True, spc=1.15, ls=1.0)
for i, (col, val, verdict, note, c) in enumerate([
    ("order_items.line_total", "$7,782,964.89", "CORRECT", "the canonical revenue measure", TEAL),
    ("invoices.amount", "$2,030,281.53", "WRONG", "invoiced to date — a different question", ERR),
]):
    yy = BODY_TOP + 0.68 + i * 0.80
    rect(s, M + 0.30, yy, tw - 0.60, 0.66, fill=WHITE, line=BORDER)
    rect(s, M + 0.30, yy, 0.05, 0.66, fill=c, line=None)
    text(s, M + 0.54, yy + 0.10, 2.90, 0.22, col, size=10.2, color=TEXT, bold=True,
         font=MONO, ls=1.0)
    text(s, M + 0.54, yy + 0.36, 2.90, 0.20, note, size=8.6, color=MUTED, ls=1.0)
    text(s, M + 3.60, yy + 0.16, 2.10, 0.30, val, size=14, color=TEXT, bold=True, ls=1.0)
    chip(s, M + tw - 1.44, yy + 0.19, 1.04, 0.28, verdict, color=c, size=8.2)
text(s, M + 0.30, BODY_TOP + 2.22, tw - 0.60, 0.22,
     "A 3.8× error — well-formed SQL, confidently returned, silently wrong.",
     size=9.6, color=WARN, bold=True, ls=1.0)

fx = M + tw + 0.36
fw = CW - tw - 0.36
card(s, fx, BODY_TOP, fw, thh, "How the graph settles it",
     [[("13 canonical metrics are stored as ", {}),
       (":Metric", {"font": MONO, "color": ACCENT_TX, "bold": True}),
       (" nodes holding the exact SQL expression and the tables it belongs to.", {})],
      [("Retrieval attaches any metric whose tables intersect the retrieved set, and the system "
        "prompt instructs the model to use the expression verbatim — never a similar-looking column.",
        {})]],
     accent=ACCENT, bsize=10.2)

yb = BODY_TOP + thh + 0.24
rect(s, M, yb, CW, 2.16, fill=SURFACE, line=BORDER)
text(s, M + 0.30, yb + 0.22, CW - 0.60, 0.26,
     "THE METRIC GLOSSARY  ·  GOVERNED, VERSIONED, SINGLE-SOURCED",
     size=8.6, color=ACCENT_TX, bold=True, spc=1.15, ls=1.0)
mets = [("total revenue", "SUM(order_items.line_total)"),
        ("average order value", "SUM(order_items.line_total) / COUNT(DISTINCT order_items.order_id)"),
        ("gross margin", "SUM((order_items.unit_price - products.cost_price) * order_items.quantity)"),
        ("target attainment rate", "SUM(sales_targets.achieved_amount) / SUM(sales_targets.target_amount)")]
for i, (n, expr) in enumerate(mets):
    ry = yb + 0.60 + i * 0.32
    text(s, M + 0.30, ry, 2.30, 0.22, n, size=9.8, color=TEXT, bold=True, ls=1.0)
    text(s, M + 2.72, ry + 0.01, CW - 3.10, 0.22, expr, size=9.4, color=SQL, font=MONO, ls=1.0)
rect(s, M + 0.30, yb + 1.84, CW - 0.60, 0.006, fill=BORDER, line=None)
text(s, M + 0.30, yb + 1.92, CW - 0.60, 0.22,
     "Defined once in schema.json, enforced everywhere. Changing a business definition is a "
     "catalogue edit — not a prompt rewrite.", size=9.2, color=MUTED, ls=1.0)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 14 — Optimisation 3: joins
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Optimisation 03  ·  Joins", "Join discovery becomes a graph traversal, not a guess",
       "The one place where using a graph database is not a preference but the correct data structure.")

for i, (t, sub, body, c) in enumerate([
    ("Without a graph", "THE MODEL INFERS",
     "It sees column names and reasons by resemblance. customer_id in two tables probably joins. "
     "Sometimes it is right. When it is wrong, the SQL still runs and returns a number that "
     "looks reasonable — the most dangerous failure mode there is.", ERR),
    ("With the graph", "THE SYSTEM RESOLVES",
     "shortestPath over :REFERENCES between the seed tables, bounded to 3 hops. Returns the "
     "actual join columns, stored from the foreign keys the catalogue declares. Nothing is "
     "inferred, so nothing can be invented.", TEAL),
]):
    x, w = cols(2, gap=0.36)[i]
    rect(s, x, BODY_TOP, w, 2.12, fill=SURFACE, line=BORDER)
    rect(s, x, BODY_TOP, w, 0.05, fill=c, line=None)
    text(s, x + 0.30, BODY_TOP + 0.26, w - 0.60, 0.22, sub, size=8.4, color=c, bold=True,
         spc=1.15, ls=1.0)
    text(s, x + 0.30, BODY_TOP + 0.54, w - 0.60, 0.30, t, size=14.5, color=TEXT, bold=True, ls=1.0)
    text(s, x + 0.30, BODY_TOP + 0.98, w - 0.60, 1.06, body, size=10.2, color=DIM, ls=1.30)

yb = BODY_TOP + 2.34
rect(s, M, yb, CW, 1.30, fill=SURFACE2, line=BORDER)
text(s, M + 0.30, yb + 0.20, CW - 0.60, 0.24,
     "A WORKED TRAVERSAL  ·  “revenue by customer segment” seeds two tables three hops apart",
     size=8.6, color=ACCENT_TX, bold=True, spc=1.1, ls=1.0)
hops = ["customers", "orders", "order_items"]
keys = ["orders.customer_id  =  customers.customer_id", "order_items.order_id  =  orders.order_id"]
CHIP_W = 1.62
TRACK_X = M + 0.34
TRACK_W = CW - 0.68
ZONE_W = (TRACK_W - 3 * CHIP_W) / 2
hx = TRACK_X
for i, h in enumerate(hops):
    chip(s, hx, yb + 0.60, CHIP_W, 0.40, h, color=SQL, size=9.6, mono=True)
    hx += CHIP_W
    if i < 2:
        arrow(s, hx + 0.06, yb + 0.73, 0.20, 0.14, color=ACCENT)
        text(s, hx + 0.32, yb + 0.66, ZONE_W - 0.42, 0.28, keys[i],
             size=8.4, color=MUTED, font=MONO, ls=1.10)
        hx += ZONE_W
text(s, TRACK_X, yb + 1.06, TRACK_W, 0.20,
     "The model receives these two conditions as fact, in the prompt, before it writes any SQL.",
     size=8.8, color=DIM, ls=1.0)

yc = yb + 1.52
for i, (t, b, c) in enumerate([
    ("Fan-out warnings", "Edges carry cardinality. On a many-to-one join the prompt warns that "
     "aggregating will double-count, and suggests COUNT(DISTINCT) or pre-aggregation.", WARN),
    ("Self-joins surfaced", "employees.manager_id → employees.employee_id and the category "
     "hierarchy are found explicitly — pairwise paths between distinct tables would miss them.", KG),
    ("Honest refusal", "When retrieved tables span domains with no join key, the context says so "
     "and instructs the model to answer separately rather than fabricate a link.", TEAL),
]):
    x, w = cols(3)[i]
    rect(s, x, yc, w, 1.08, fill=SURFACE, line=BORDER)
    rect(s, x, yc, 0.05, 1.08, fill=c, line=None)
    text(s, x + 0.26, yc + 0.16, w - 0.50, 0.24, t, size=11, color=TEXT, bold=True, ls=1.0)
    text(s, x + 0.26, yc + 0.46, w - 0.50, 0.52, b, size=9.2, color=DIM, ls=1.24)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 15 — Trust & governance
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Trust & Governance", "An answer a leader can act on has to show its work",
       "Four independent mechanisms, none of which depends on trusting the model.")

mech = [
    ("01", "The reasoning trace", "Every tool call is recorded as a Step — which tables were "
     "discovered, what SQL was attempted, where the agent corrected itself. Returned with every "
     "response and rendered live in the UI.", ACCENT),
    ("02", "Deterministic grounding", "After the answer is produced, the figures it states are "
     "matched against the returned rows. Zero extra model calls, fully deterministic.", TEAL),
    ("03", "Honest success", "The success flag is true only when the run reached a conclusive "
     "answer that actually executed SQL. A step-exhausted run reports itself as unresolved.", SQL),
    ("04", "Graceful failure", "Neo4j paused, Groq rate-limited or credentials missing return a "
     "classified, human-readable envelope. A stack trace never reaches the user.", WARN),
]
for i, (n, t, b, c) in enumerate(mech):
    x, w = cols(4)[i]
    rect(s, x, BODY_TOP, w, 2.42, fill=SURFACE, line=BORDER)
    rect(s, x, BODY_TOP, w, 0.05, fill=c, line=None)
    text(s, x + 0.26, BODY_TOP + 0.28, 0.50, 0.28, n, size=15, color=c, bold=True, font=LIGHT, ls=1.0)
    text(s, x + 0.26, BODY_TOP + 0.66, w - 0.50, 0.30, t, size=12.4, color=TEXT, bold=True, ls=1.05)
    text(s, x + 0.26, BODY_TOP + 1.08, w - 0.50, 1.16, b, size=9.8, color=DIM, ls=1.28)

yb = BODY_TOP + 2.64
rect(s, M, yb, 7.55, 1.80, fill=SURFACE2, line=TEAL, lw=1.4)
rect(s, M, yb, 0.055, 1.80, fill=TEAL, line=None)
text(s, M + 0.30, yb + 0.22, 7.00, 0.26, "GROUNDING  ·  ADVISORY BY DESIGN, NOT A GATE",
     size=8.6, color=TEAL, bold=True, spc=1.15, ls=1.0)
text(s, M + 0.30, yb + 0.56, 7.00, 1.04,
     [[("The check is deliberately never allowed to block an answer. It false-negatives on derived "
        "figures a correct answer computes but no single cell contains — “4.06 orders per customer”, "
        "for instance. ", {"color": DIM, "size": 10.4}),
       ("So it is surfaced as an “unverified against data” caveat, and the eval harness measures its "
        "false-negative rate rather than trusting it blindly.",
        {"color": TEXT, "size": 10.4, "bold": True})]], ls=1.30)

ox3 = M + 7.55 + 0.36
ow3 = CW - 7.55 - 0.36
rect(s, ox3, yb, ow3, 1.80, fill=SURFACE, line=BORDER)
text(s, ox3 + 0.28, yb + 0.22, ow3 - 0.56, 0.26, "WHAT COMES BACK, EVERY TIME", size=8.6,
     color=ACCENT_TX, bold=True, spc=1.15, ls=1.0)
for i, (k, v) in enumerate([("answer", "plain-language response"),
                            ("sql", "the exact query executed"),
                            ("rows / columns", "the data behind it"),
                            ("trace[]", "every reasoning step"),
                            ("grounded", "verified · or caveated")]):
    ry = yb + 0.56 + i * 0.24
    text(s, ox3 + 0.28, ry, 1.60, 0.20, k, size=8.8, color=SQL, bold=True, font=MONO, ls=1.0)
    text(s, ox3 + 1.96, ry, ow3 - 2.24, 0.20, v, size=8.8, color=DIM, ls=1.0)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 16 — Measurable accuracy
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Measurable Accuracy", "“Accurate” is only a claim until it is falsifiable",
       "An offline evaluation harness makes answer quality a number we can regression-test.")

for i, (x, w) in enumerate(cols(4)):
    v, l, nt, c = [
        ("15", "Gold questions", "Committed expected values across all five domains, including every known trap.", ACCENT_TX),
        ("184", "Automated tests", "The loop, tools, retriever, scoring and grounding — all covered.", SQL),
        ("0", "Live calls needed", "A scripted fake model runs the harness deterministically in CI.", TEAL),
        ("42", "Fixed data seed", "A determinism gate fails loudly if generated data ever drifts.", WARN),
    ][i]
    kpi(s, x, BODY_TOP, w, 1.74, v, l, nt, c)

yb = BODY_TOP + 1.98
card(s, M, yb, 5.95, 2.46, "How the ruler works",
     [[("result_match", {"font": MONO, "color": ACCENT_TX, "bold": True}),
       (" scores values as a tolerant multiset: it ignores column aliases, row order, "
        "INTEGER-vs-REAL and formatting.", {})],
      [("Right answer, wrong shape counts as a match — a ruler that failed those would "
        "mis-rank every downstream fix.", {"color": TEXT, "bold": True})],
      [("Wrong answer still fails: the revenue trap is rejected outright. Adversarial pairs are "
        "themselves unit-tested.", {})]],
     accent=ACCENT, bsize=10.2)

card(s, M + 6.31, yb, CW - 6.31, 2.46, "What the harness protects against",
     None, accent=TEAL)
prot = [("Silent data drift", "Every gold SQL is re-run against the tracked database on each test run."),
        ("Prompt regressions", "Prompt changes on a 70B model can hurt as easily as help — the harness A/Bs them."),
        ("Over-claiming", "Cross-domain questions are scored on graceful refusal, not on producing a number."),
        ("Metric substitution", "The revenue trap is a permanent, committed test case.")]
for i, (t, b) in enumerate(prot):
    ry = yb + 0.62 + i * 0.46
    text(s, M + 6.59, ry, CW - 6.87, 0.22, t, size=10.4, color=TEXT, bold=True, ls=1.0)
    text(s, M + 6.59, ry + 0.21, CW - 6.87, 0.24, b, size=9.0, color=DIM, ls=1.0)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 17 — Worked example
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Worked Example", "One question, end to end",
       "Exactly what the user sees, and exactly what the system did to earn it.")

rect(s, M, BODY_TOP, CW, 0.66, fill=SURFACE2, line=ACCENT, lw=1.4)
rect(s, M, BODY_TOP, 0.055, 0.66, fill=ACCENT, line=None)
text(s, M + 0.30, BODY_TOP + 0.13, 1.00, 0.24, "ASKED", size=8.4, color=ACCENT_TX, bold=True,
     spc=1.15, ls=1.0)
text(s, M + 1.34, BODY_TOP + 0.15, CW - 1.70, 0.34,
     "“What was total revenue by region last quarter?”", size=15, color=TEXT, bold=True, ls=1.0)

trace = [
    ("STEP 1", "get_schema_context", "Routes to the Sales domain. Seeds order_items, orders, "
     "customers, regions. Walks :REFERENCES for the join keys. Attaches the total revenue metric.",
     "4 tables · 3 joins · 1 metric", KG),
    ("STEP 2", "sample_values", "Reads the distinct values of regions.region_name so the filter "
     "uses a real literal rather than a guessed one.", "6 values returned", SQL),
    ("STEP 3", "run_sql", "Executes the SELECT with the graph-supplied join keys and the canonical "
     "revenue expression, read-only and row-capped.", "6 rows · 4 columns", SQL),
    ("STEP 4", "final answer", "States the figures in plain language, then the grounding check "
     "confirms each one appears in the returned rows.", "grounded = true", TEAL),
]
ty2 = BODY_TOP + 0.88
for i, (n, tool, b, obs, c) in enumerate(trace):
    yy = ty2 + i * 0.80
    rect(s, M, yy, 8.55, 0.70, fill=SURFACE, line=BORDER)
    rect(s, M, yy, 0.05, 0.70, fill=c, line=None)
    text(s, M + 0.24, yy + 0.13, 0.70, 0.20, n, size=7.8, color=MUTED, bold=True, spc=0.9, ls=1.0)
    text(s, M + 0.24, yy + 0.36, 1.90, 0.22, tool, size=9.8, color=c, bold=True, font=MONO, ls=1.0)
    text(s, M + 2.32, yy + 0.14, 4.20, 0.46, b, size=9.2, color=DIM, ls=1.20)
    chip(s, M + 6.68, yy + 0.21, 1.74, 0.28, obs, color=c, size=7.8)

ax2 = M + 8.55 + 0.36
aw2 = CW - 8.55 - 0.36
rect(s, ax2, ty2, aw2, 3.26, fill=SURFACE, line=TEAL, lw=1.4)
text(s, ax2 + 0.28, ty2 + 0.22, aw2 - 0.56, 0.26, "DELIVERED TO THE USER", size=8.6,
     color=TEAL, bold=True, spc=1.15, ls=1.0)
for i, (k, v, c) in enumerate([
    ("Answer", "Plain language, stating each figure.", TEXT),
    ("SQL", "The exact query executed.", SQL),
    ("Table", "The 6 returned rows, inline.", SQL),
    ("Trace", "All four steps, expandable.", KG),
    ("Badge", "Verified against data.", TEAL),
]):
    ry = ty2 + 0.56 + i * 0.55
    text(s, ax2 + 0.28, ry, aw2 - 0.56, 0.22, k, size=10, color=c, bold=True, ls=1.0)
    text(s, ax2 + 0.28, ry + 0.21, aw2 - 0.56, 0.22, v, size=8.8, color=DIM, ls=1.18)
    if i < 4:
        rect(s, ax2 + 0.28, ry + 0.46, aw2 - 0.56, 0.006, fill=BORDER, line=None)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 18 — Deployment & operations
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Deployment & Operations", "Running today, on managed infrastructure",
       "One process serves the API and the UI. Every push to the branch ships automatically.")

stack = [("Application", "FastAPI + Uvicorn", "One process serves /api/* and the UI from a single origin — no CORS, no second deployment.", ACCENT),
         ("Knowledge graph", "Neo4j AuraDB", "Managed cloud. Vector indexes and the schema graph. Rebuilt only when the catalogue changes.", KG),
         ("Language model", "LLaMA-3.3-70B on Groq", "Fast tool-calling inference. The provider sits behind one callable and is swappable.", SQL),
         ("Embeddings", "all-MiniLM-L6-v2, local", "Runs in the container. No per-query API cost, no external dependency.", TEAL),
         ("Execution store", "SQLite, tracked in-repo", "Zero infrastructure. Opened read-only for every query.", WARN),
         ("Hosting", "Railway, auto-deploy", "Health-checked releases from GitHub. Configuration is committed as railway.toml.", ACCENT_TX)]
for i, (role, tech, note, c) in enumerate(stack):
    x, w = cols(3)[i % 3]
    yy = BODY_TOP + (i // 3) * 1.30
    rect(s, x, yy, w, 1.18, fill=SURFACE, line=BORDER)
    rect(s, x, yy, 0.05, 1.18, fill=c, line=None)
    text(s, x + 0.26, yy + 0.16, w - 0.50, 0.20, role.upper(), size=8, color=MUTED, bold=True,
         spc=1.0, ls=1.0)
    text(s, x + 0.26, yy + 0.40, w - 0.50, 0.24, tech, size=11.6, color=c, bold=True, ls=1.0)
    text(s, x + 0.26, yy + 0.70, w - 0.50, 0.40, note, size=9.2, color=DIM, ls=1.22)

yb = BODY_TOP + 2.72
rect(s, M, yb, CW, 1.72, fill=SURFACE2, line=BORDER)
text(s, M + 0.30, yb + 0.22, CW - 0.60, 0.26, "OPERATIONAL POSTURE", size=8.6,
     color=ACCENT_TX, bold=True, spc=1.15, ls=1.0)
ops = [("Health endpoint", "/api/health reports graph reachability and catalogue freshness, with a bounded 5-second probe so a paused instance never hangs the poll."),
       ("Staleness detection", "A build stamp in the graph is compared against the catalogue fingerprint — a stale graph is reported, never silently served."),
       ("Rate-limit discipline", "Transient limits honour Retry-After; a sustained quota raises a distinct signal so callers stop rather than hammer."),
       ("Config as code", "Start command, health check and restart policy live in railway.toml — the deployment is reviewable in the repository.")]
OPS_COL = [(M + 0.30, 5.75), (M + 6.35, 5.75)]
for i, (t, b) in enumerate(ops):
    x, w = OPS_COL[i % 2]
    ry = yb + 0.58 + (i // 2) * 0.56
    text(s, x, ry, 1.90, 0.22, t, size=10, color=TEXT, bold=True, ls=1.0)
    text(s, x + 1.98, ry, w - 2.02, 0.46, b, size=9.0, color=DIM, ls=1.18)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 19 — Business impact + roadmap
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Business Impact", "What changes for the organisation",
       "The technical work matters because of what it removes: a queue, a dependency, and a class of silent error.")

imp = [("From days to seconds", "A question that used to become a ticket is now asked and answered in the same conversation.", ACCENT),
       ("Analysts move up the stack", "The team stops writing routine joins and starts curating the catalogue and the metric glossary.", KG),
       ("One definition of the truth", "“Revenue” means one expression, governed in one file — the same number in every answer, for everyone.", SQL),
       ("Auditable by default", "Every answer arrives with its SQL, its rows and its reasoning. Nothing has to be taken on trust.", TEAL)]
for i, (t, b, c) in enumerate(imp):
    x, w = cols(4)[i]
    rect(s, x, BODY_TOP, w, 2.14, fill=SURFACE, line=BORDER)
    rect(s, x, BODY_TOP, w, 0.05, fill=c, line=None)
    text(s, x + 0.28, BODY_TOP + 0.34, w - 0.56, 0.60, t, size=15, color=TEXT, bold=True, ls=1.10)
    text(s, x + 0.28, BODY_TOP + 1.06, w - 0.56, 0.90, b, size=10, color=DIM, ls=1.30)

yb = BODY_TOP + 2.38
eyebrow(s, M, yb, "Where this goes next", color=MUTED, w=6.0)
road = [("NEAR", "Widen the catalogue", "Point the same pipeline at real warehouse metadata. The architecture is catalogue-driven — new domains need no code change.", ACCENT),
        ("NEXT", "Raise the accuracy floor", "Run the live evaluation baseline, then use it to gate prompt hardening and a higher step budget.", SQL),
        ("LATER", "Specialise the agents", "Split planner, SQL author and synthesiser so each can be tuned and evaluated independently.", KG)]
for i, (tag, t, b, c) in enumerate(road):
    x, w = cols(3)[i]
    yy = yb + 0.30
    rect(s, x, yy, w, 1.66, fill=SURFACE2, line=BORDER)
    rect(s, x, yy, 0.05, 1.66, fill=c, line=None)
    chip(s, x + 0.26, yy + 0.20, 0.90, 0.26, tag, color=c, size=7.6)
    text(s, x + 0.26, yy + 0.60, w - 0.50, 0.26, t, size=12, color=TEXT, bold=True, ls=1.0)
    text(s, x + 0.26, yy + 0.92, w - 0.50, 0.62, b, size=9.4, color=DIM, ls=1.26)
footer(s)

# ═══════════════════════════════════════════════════════════════════════════
# 20 — Close
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(wash=True)
rect(s, 0, 0, SW, 0.14, fill=ACCENT, line=None)
rect(s, M, 1.30, 0.055, 0.34, fill=ACCENT, line=None)
eyebrow(s, M + 0.17, 1.34, "In one line", w=6.0)

text(s, M, 1.86, 11.6, 1.72,
     [[("The knowledge graph is what makes", {"color": TEXT})],
      [("the language model ", {"color": TEXT}),
       ("trustworthy.", {"color": ACCENT_TX, "bold": True})]],
     size=38, ls=1.18, font=LIGHT)

rect(s, M, 3.74, 5.0, 0.014, fill=ACCENT, line=None)
text(s, M, 3.98, 11.6, 0.94,
     [[("The model supplies language. The graph supplies structure — which tables are relevant, "
        "how they join, and what the business means by a metric. ", {"color": DIM})],
      [("Separating those two responsibilities is the whole design, and it is why the answers hold up.",
        {"color": TEXT, "bold": True})]],
     size=13.5, ls=1.44)

for i, (x, w) in enumerate(cols(4, gap=0.22)):
    v, l = [("83.8%", "less prompt context"), ("6.2×", "smaller payload"),
            ("441", "graph nodes serving the plan"), ("184", "tests guarding the loop")][i]
    rect(s, x, 5.34, w, 0.98, fill=SURFACE, line=BORDER)
    rect(s, x, 5.34, w, 0.05, fill=ACCENT, line=None)
    text(s, x, 5.56, w, 0.34, v, size=20, color=TEXT, bold=True, align=PP_ALIGN.CENTER, ls=1.0)
    text(s, x, 5.96, w, 0.22, l.upper(), size=7.8, color=ACCENT_TX, bold=True,
         align=PP_ALIGN.CENTER, spc=0.9, ls=1.0)

text(s, M, 6.72, 11.6, 0.24,
     "DataPulse   ·   Knowledge-Graph Grounded Text-to-SQL   ·   Live on Railway",
     size=9.5, color=MUTED, spc=0.4)

prs.save(OUT)
print("saved:", OUT)
