from __future__ import annotations
import math
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = r"c:\Data_pulse_kg\DataPulse_Leadership_Deck.pptx"

BG=RGBColor(0xFF,0xFF,0xFF); PAPER=RGBColor(0xFA,0xFA,0xFC); SURFACE=RGBColor(0xF6,0xF6,0xF9)
SURFACE2=RGBColor(0xEF,0xEF,0xF4); BORDER=RGBColor(0xE1,0xE1,0xE9); BORDER_HI=RGBColor(0xC6,0xC6,0xD2)
ACCENT=RGBColor(0xA1,0x00,0xFF); ACCENT_TX=RGBColor(0x6D,0x00,0xB3); ACCENT_DK=RGBColor(0x4A,0x00,0x7A)
KG=RGBColor(0x6D,0x00,0xB3); SQL=RGBColor(0x0A,0x6B,0xAF); TEXT=RGBColor(0x11,0x11,0x19)
DIM=RGBColor(0x45,0x45,0x4F); MUTED=RGBColor(0x78,0x78,0x8A); TEAL=RGBColor(0x0B,0x7C,0x63)
WHITE=RGBColor(0xFF,0xFF,0xFF); T1=RGBColor(0xEE,0xDD,0xFB); T2=RGBColor(0xD5,0xB0,0xF6)
T3=RGBColor(0xB0,0x6E,0xEC); T4=RGBColor(0x8A,0x28,0xDE)
SANS="Segoe UI"; LIGHT="Segoe UI Light"; MONO="Consolas"
SW,SH=13.333,7.5; M=0.60; CW=SW-2*M; CXC=M+CW/2; FOOT_Y=6.99

prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
BLANK=prs.slide_layouts[6]; _n=[0]

def _style(shape,fill=None,line=None,lw=1.0):
    shape.shadow.inherit=False
    if fill is None: shape.fill.background()
    else: shape.fill.solid(); shape.fill.fore_color.rgb=fill
    if line is None: shape.line.fill.background()
    else: shape.line.color.rgb=line; shape.line.width=Pt(lw)
    if shape.has_text_frame: shape.text_frame.word_wrap=True
    return shape

def new_slide(fill=BG):
    s=prs.slides.add_slide(BLANK)
    _style(s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(SW),Inches(SH)),fill=fill,line=None)
    return s

def rect(s,x,y,w,h,fill=SURFACE,line=BORDER,lw=1.0,shape=MSO_SHAPE.RECTANGLE,rot=None):
    sh=_style(s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h)),fill,line,lw)
    if rot is not None: sh.rotation=rot
    return sh

def grad(s,x,y,w,h,c1,c2,angle=0.0,shape=MSO_SHAPE.RECTANGLE):
    sh=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.shadow.inherit=False; sh.line.fill.background()
    try:
        f=sh.fill; f.gradient()
        f.gradient_stops[0].color.rgb=c1; f.gradient_stops[0].position=0.0
        f.gradient_stops[1].color.rgb=c2; f.gradient_stops[1].position=1.0
        f.gradient_angle=angle
    except Exception: sh.fill.solid(); sh.fill.fore_color.rgb=c1
    if sh.has_text_frame: sh.text_frame.word_wrap=True
    return sh

def tracking(run,pts): run.font._rPr.set("spc",str(int(pts*100)))

def text(s,x,y,w,h,spans,size=10.5,color=DIM,bold=False,font=SANS,
         align=PP_ALIGN.LEFT,ls=1.24,anchor=MSO_ANCHOR.TOP,spc=0.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    paras=[spans] if isinstance(spans,str) else spans
    for i,para in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=ls
        runs=[(para,{})] if isinstance(para,str) else para
        runs=[(r,{}) if isinstance(r,str) else r for r in runs]
        for txt,ov in runs:
            r=p.add_run(); r.text=txt; f=r.font
            f.name=ov.get("font",font); f.size=Pt(ov.get("size",size))
            f.bold=ov.get("bold",bold); f.color.rgb=ov.get("color",color)
            t=ov.get("spc",spc)
            if t: tracking(r,t)
    return tb

def ctext(s,x,y,w,h,spans,**kw):
    kw.setdefault("align",PP_ALIGN.CENTER); return text(s,x,y,w,h,spans,**kw)

def label_in(s,x,y,w,h,spans,size=10.5,color=TEXT,bold=True,font=SANS,ls=1.10):
    return text(s,x,y,w,h,spans,size=size,color=color,bold=bold,font=font,
                align=PP_ALIGN.CENTER,ls=ls,anchor=MSO_ANCHOR.MIDDLE)

def eyebrow(s,x,y,label,color=ACCENT_TX,w=None,size=8.5):
    text(s,x,y,w or 7.0,0.20,label.upper(),size=size,color=color,bold=True,spc=1.3)

def header(s,kicker,title):
    rect(s,M,0.44,0.055,0.30,fill=ACCENT,line=None)
    eyebrow(s,M+0.17,0.47,kicker)
    text(s,M,0.70,CW,0.46,title,size=25,color=TEXT,bold=True,ls=1.0)
    rect(s,M,1.34,CW,0.014,fill=ACCENT,line=None)
    return 1.34

def footer(s,label="DataPulse  ·  Knowledge-Graph Grounded Text-to-SQL"):
    _n[0]+=1
    text(s,M,FOOT_Y,8.0,0.20,label,size=8,color=MUTED,spc=0.4)
    text(s,SW-M-2.0,FOOT_Y,2.0,0.20,f"{_n[0]:02d}",size=8,color=ACCENT_TX,
         align=PP_ALIGN.RIGHT,bold=True)

def cols(n,gap=0.24,x0=M,total=CW):
    w=(total-gap*(n-1))/n; return [(x0+i*(w+gap),w) for i in range(n)]

def arrow(s,x,y,w,h,shape=MSO_SHAPE.RIGHT_ARROW,color=BORDER_HI,rot=None):
    return rect(s,x,y,w,h,fill=color,line=None,shape=shape,rot=rot)

def line(s,x1,y1,x2,y2,color=BORDER_HI,w=1.25,dash=False):
    ln=s.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    ln.line.color.rgb=color; ln.line.width=Pt(w)
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        ln.line.dash_style=MSO_LINE_DASH_STYLE.DASH
    return ln

def node(s,x,y,w,h,title,sub=None,color=KG,fill=WHITE,tsize=10.6,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE,mono=False):
    rect(s,x,y,w,h,fill=fill,line=color,lw=1.5,shape=shape)
    if sub:
        text(s,x,y+h/2-0.30,w,0.24,title,size=tsize,color=color,bold=True,
             font=MONO if mono else SANS,align=PP_ALIGN.CENTER,ls=1.0)
        text(s,x,y+h/2+0.00,w,0.20,sub,size=7.8,color=MUTED,align=PP_ALIGN.CENTER,ls=1.0)
    else:
        label_in(s,x,y,w,h,title,size=tsize,color=color,font=MONO if mono else SANS)
    return (x+w/2,y+h/2)

def hero(s,x,y,w,value,label,color=ACCENT_TX,vsize=54,lsize=9):
    ctext(s,x,y,w,0.80,value,size=vsize,color=color,bold=True,font=LIGHT,ls=0.98)
    ctext(s,x,y+0.86,w,0.24,label.upper(),size=lsize,color=MUTED,bold=True,spc=1.2,ls=1.0)

def tile(s,x,y,w,h,value,label,color=ACCENT_TX,vsize=20):
    rect(s,x,y,w,h,fill=SURFACE,line=BORDER)
    rect(s,x,y,w,0.05,fill=color,line=None)
    ctext(s,x,y+h*0.26,w,0.34,value,size=vsize,color=TEXT,bold=True,ls=1.0)
    ctext(s,x,y+h*0.66,w,0.22,label.upper(),size=7.8,color=color,bold=True,spc=0.9,ls=1.0)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title (full-bleed gradient, no footer)
# ─────────────────────────────────────────────────────────────────────────────
s1 = new_slide(fill=ACCENT_DK)
grad(s1, 0, 0, SW, SH, ACCENT_DK, ACCENT, angle=315)

ctext(s1, 0, 2.0, SW, 1.0, "DataPulse", size=62, color=WHITE, bold=True, font=LIGHT, ls=1.0)
ctext(s1, 0, 3.1, SW, 0.40, "Knowledge Graph · Text-to-SQL · Production",
      size=16, color=WHITE, bold=False, ls=1.0)
ctext(s1, 0, 3.7, SW, 0.30, "Accenture  |  August 2026",
      size=11, color=T2, bold=False, ls=1.0)

# Three badge pills at y=4.4
badges = ["Neo4j AuraDB", "LLaMA-3.3-70b · Groq", "FastAPI · SQLite"]
bw = 2.6; gap_b = 0.40; total_b = len(badges)*bw + (len(badges)-1)*gap_b
bx_start = (SW - total_b) / 2
for i, badge_text in enumerate(badges):
    bx = bx_start + i*(bw+gap_b)
    rect(s1, bx, 4.4, bw, 0.44, fill=WHITE, line=ACCENT, lw=1.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    label_in(s1, bx, 4.4, bw, 0.44, badge_text, size=10, color=ACCENT_TX)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
s2 = new_slide()
header(s2, "context", "Enterprise data is locked in complexity")

card_cols = cols(3)
card_y = 1.55; card_h = 4.60
card_colors = [ACCENT, SQL, TEAL]
card_titles = ["Schema Overload", "Hallucination Risk", "Context Waste"]
card_bodies = [
    "50 tables · 23k rows\nno LLM knows full layout",
    "LLMs guess column names,\nfabricate joins, invent metrics",
    "Entire schema in every prompt\nexhausts token budget",
]
stat_vals  = ["10,224 tokens", "~40%", "6.2×"]
stat_labs  = ["full-schema prompt", "error rate without grounding", "tokens wasted vs KG slice"]

for i, (cx, cw) in enumerate(card_cols):
    rect(s2, cx, card_y, cw, card_h, fill=SURFACE, line=BORDER)
    rect(s2, cx, card_y, cw, 0.06, fill=card_colors[i], line=None)
    text(s2, cx+0.12, card_y+0.14, cw-0.24, 0.34, card_titles[i],
         size=13, color=TEXT, bold=True, ls=1.0)
    text(s2, cx+0.12, card_y+0.56, cw-0.24, 0.80, card_bodies[i],
         size=10, color=DIM, ls=1.3)
    # stat tile at bottom of card
    tile_y = card_y + card_h - 1.10
    tile(s2, cx+0.10, tile_y, cw-0.20, 0.90,
         stat_vals[i], stat_labs[i], color=card_colors[i], vsize=16)

footer(s2)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Solution in One Picture
# ─────────────────────────────────────────────────────────────────────────────
s3 = new_slide()
header(s3, "solution", "DataPulse: ask in plain English, get SQL-backed answers")

# Five boxes + arrows in a horizontal flow
flow_y = 2.85; flow_h = 0.80
arrow_w = 0.35; arrow_h = 0.28

boxes = [
    # (w, fill, line, text_color, label)
    (2.0, SURFACE2, BORDER_HI, TEXT,  "Natural Language\nQuestion"),
    (2.0, T1,       KG,        KG,    "Knowledge Graph\nRetriever"),
    (2.1, ACCENT,   None,      WHITE, "ReAct Agent\nLLaMA-3.3-70b"),
    (2.0, RGBColor(0xE3,0xF2,0xFD), SQL, SQL, "SQLite\nQuery"),
    (2.0, SURFACE2, TEAL,      TEAL,  "Verified\nAnswer"),
]

# Calculate total width to centre the flow
total_flow_w = sum(b[0] for b in boxes) + arrow_w * (len(boxes)-1)
fx = (SW - total_flow_w) / 2

box_positions = []
cx = fx
for i, (bw, bf, bl, btc, bt) in enumerate(boxes):
    bl_arg = bl if bl is not None else None
    rect(s3, cx, flow_y, bw, flow_h,
         fill=bf, line=bl_arg if bl_arg else bf, lw=1.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # multiline label
    lines_t = bt.split("\n")
    if len(lines_t) == 2:
        text(s3, cx, flow_y+0.14, bw, 0.26, lines_t[0],
             size=9.5, color=btc, bold=True, align=PP_ALIGN.CENTER, ls=1.0)
        text(s3, cx, flow_y+0.42, bw, 0.26, lines_t[1],
             size=9.0, color=btc, bold=False, align=PP_ALIGN.CENTER, ls=1.0)
    else:
        label_in(s3, cx, flow_y, bw, flow_h, bt, size=9.5, color=btc)
    box_positions.append((cx, bw))
    cx += bw
    if i < len(boxes) - 1:
        arrow(s3, cx, flow_y + (flow_h - arrow_h)/2, arrow_w, arrow_h,
              shape=MSO_SHAPE.RIGHT_ARROW, color=BORDER_HI)
        cx += arrow_w

# Sub-labels below specific boxes
sub_labels = [
    (1, "Neo4j · 441 nodes · 384-d vectors"),
    (2, "≤6 steps · tool-calling loop"),
    (3, "50 tables · read-only"),
]
for idx, lbl in sub_labels:
    bx, bw2 = box_positions[idx]
    text(s3, bx, flow_y + flow_h + 0.08, bw2, 0.22, lbl,
         size=7.8, color=MUTED, align=PP_ALIGN.CENTER, ls=1.0)

# Response time badge bottom-right
rect(s3, 9.8, 5.90, 3.10, 0.44, fill=SURFACE2, line=BORDER_HI,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
label_in(s3, 9.8, 5.90, 3.10, 0.44,
         "3–8 s end-to-end · Groq LPU", size=9, color=ACCENT_TX)

footer(s3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — End-to-End Architecture
# ─────────────────────────────────────────────────────────────────────────────
s4 = new_slide()
header(s4, "architecture", "Full-stack: Vercel → Render → Groq → Neo4j → SQLite")

WARN = RGBColor(0xFF, 0x6B, 0x00)

# ROW 1: User Browser
rect(s4, 0.6, 1.55, 2.2, 0.70, fill=SURFACE2, line=BORDER_HI, lw=1.5,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
label_in(s4, 0.6, 1.55, 2.2, 0.70, "User Browser", size=10, color=TEXT)

# Arrow down from browser to Vercel
line(s4, 1.7, 2.25, 1.7, 2.55, color=BORDER_HI, w=1.5)

# ROW 2: Vercel CDN
rect(s4, 0.6, 2.55, 2.2, 0.70, fill=T1, line=KG, lw=1.5,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
label_in(s4, 0.6, 2.55, 2.2, 0.70, "Vercel CDN", size=10, color=KG)

# Arrow right: Vercel → FastAPI
arrow(s4, 2.80, 2.83, 0.40, 0.24, shape=MSO_SHAPE.RIGHT_ARROW, color=BORDER_HI)

# Render / FastAPI
rect(s4, 3.20, 2.55, 2.8, 0.70, fill=ACCENT, line=None,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
label_in(s4, 3.20, 2.55, 2.8, 0.70, "Render / FastAPI", size=10, color=WHITE)

# Arrow right: FastAPI → Groq
arrow(s4, 6.00, 2.83, 0.40, 0.24, shape=MSO_SHAPE.RIGHT_ARROW, color=BORDER_HI)

# Groq / LLaMA
rect(s4, 6.40, 2.55, 2.8, 0.70, fill=WARN, line=None,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
label_in(s4, 6.40, 2.55, 2.8, 0.70, "Groq / LLaMA-3.3-70b", size=9.5, color=WHITE)

# tool calls label
text(s4, 6.40, 3.28, 2.8, 0.22, "← tool calls →",
     size=7.8, color=MUTED, align=PP_ALIGN.CENTER, ls=1.0)

# ROW 3: Neo4j + SQLite below FastAPI
# Line from FastAPI down to Neo4j
line(s4, 4.60, 3.25, 4.60, 3.55, color=KG, w=1.25)

# Neo4j AuraDB
rect(s4, 2.8, 3.55, 2.8, 0.70, fill=T2, line=KG, lw=1.5,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s4, 2.8, 3.62, 2.8, 0.26, "Neo4j AuraDB",
     size=10, color=KG, bold=True, align=PP_ALIGN.CENTER, ls=1.0)
text(s4, 2.8, 3.90, 2.8, 0.20, "441 nodes · 3 vector indexes",
     size=8, color=DIM, align=PP_ALIGN.CENTER, ls=1.0)

# Diagonal line from FastAPI to SQLite
line(s4, 5.20, 3.25, 7.60, 3.55, color=SQL, w=1.25)

# SQLite DB
rect(s4, 6.20, 3.55, 2.8, 0.70, fill=RGBColor(0xE3,0xF2,0xFD), line=SQL, lw=1.5,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s4, 6.20, 3.62, 2.8, 0.26, "SQLite",
     size=10, color=SQL, bold=True, align=PP_ALIGN.CENTER, ls=1.0)
text(s4, 6.20, 3.90, 2.8, 0.20, "50 tables · 23k rows",
     size=8, color=DIM, align=PP_ALIGN.CENTER, ls=1.0)

# ROW 4: fastembed ONNX
line(s4, 4.20, 4.25, 4.20, 4.55, color=BORDER_HI, w=1.0)
rect(s4, 2.8, 4.55, 2.8, 0.60, fill=SURFACE, line=BORDER, lw=1.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s4, 2.8, 4.64, 2.8, 0.24, "fastembed ONNX",
     size=9.5, color=DIM, bold=True, align=PP_ALIGN.CENTER, ls=1.0)
text(s4, 2.8, 4.88, 2.8, 0.20, "all-MiniLM-L6-v2 · 384-d",
     size=8, color=MUTED, align=PP_ALIGN.CENTER, ls=1.0)

# Right side legend panel
rect(s4, 10.20, 1.55, 2.85, 4.20, fill=SURFACE, line=BORDER, lw=1.0)
text(s4, 10.30, 1.62, 2.65, 0.26, "Tech Stack",
     size=11, color=TEXT, bold=True, ls=1.0)
rect(s4, 10.20, 1.88, 2.85, 0.025, fill=BORDER, line=None)
stack_items = [
    "FastAPI · uvicorn",
    "Neo4j Python driver",
    "groq SDK",
    "fastembed",
    "python-dotenv",
    "aiofiles",
    "pandas",
]
for j, item in enumerate(stack_items):
    text(s4, 10.30, 1.96 + j*0.45, 2.65, 0.38, item,
         size=9, color=DIM, ls=1.0)

footer(s4)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Knowledge Graph
# ─────────────────────────────────────────────────────────────────────────────
s5 = new_slide()
header(s5, "knowledge graph", "Schema encoded as a navigable node-link graph")

# Left half: 4-level hierarchy
lx_center = 2.5

# Level 1 — DOMAIN
node(s5, lx_center - 0.9, 2.0, 1.8, 0.60,
     "DOMAIN", sub="5 nodes", color=ACCENT, fill=T1)

# Lines down Level1 → Level2
line(s5, lx_center, 2.60, lx_center, 3.00, color=ACCENT, w=1.25)

# Level 2 — TABLE (single representative)
node(s5, lx_center - 0.9, 3.00, 1.8, 0.60,
     "TABLE", sub="50 nodes", color=KG, fill=WHITE)

# Lines down Level2 → Level3 (two columns)
line(s5, lx_center, 3.60, 1.75, 4.00, color=KG, w=1.0)
line(s5, lx_center, 3.60, 3.25, 4.00, color=KG, w=1.0)

# Level 3 — COLUMN (two shown)
node(s5, 1.0, 4.00, 1.5, 0.55, "COLUMN", sub="373 nodes", color=SQL, fill=RGBColor(0xE3,0xF2,0xFD))
node(s5, 2.5, 4.00, 1.5, 0.55, "COLUMN", sub="373 nodes", color=SQL, fill=RGBColor(0xE3,0xF2,0xFD))

# Line down Level3 → Level4
line(s5, lx_center, 4.55, lx_center, 5.00, color=SQL, w=1.0)

# Level 4 — METRIC
node(s5, lx_center - 0.75, 5.00, 1.5, 0.55,
     "METRIC", sub="13 nodes", color=TEAL, fill=RGBColor(0xE0,0xF4,0xF0))

# Right half stats 2×2 grid
stat_cols = cols(2, gap=0.22, x0=7.0, total=5.60)
stat_rows = [
    [("441", "Total Nodes"), ("40", ":REFERENCES edges")],
    [("3", "Vector indexes"), ("384-d", "Embedding dim")],
]
for ri, row in enumerate(stat_rows):
    for ci, (val, lbl) in enumerate(row):
        sx, sw2 = stat_cols[ci]
        sy = 2.10 + ri * 1.55
        tile(s5, sx, sy, sw2, 1.30, val, lbl, color=ACCENT_TX, vsize=26)

# Bottom query note
text(s5, M, 6.55, CW, 0.26,
     "Query: embed question → cosine search → shortestPath join",
     size=8.5, color=MUTED, ls=1.0)

footer(s5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Context Economy
# ─────────────────────────────────────────────────────────────────────────────
s6 = new_slide()
header(s6, "context economy", "Knowledge graph cuts LLM context by 83.8%")

col_y = 1.60; max_h = 4.30

# Without KG column
wx = 0.90; ww = 4.50
text(s6, wx, col_y, ww, 0.34, "Without KG",
     size=13, color=DIM, bold=True, ls=1.0)
# Bar height proportional: 10224/10224 * 3.20 = 3.20
without_h = 3.20
rect(s6, wx + 0.80, col_y + 0.44, 2.9, without_h, fill=T1, line=BORDER_HI, lw=1.5)
ctext(s6, wx + 0.80, col_y + 0.44 + without_h/2 - 0.18, 2.9, 0.36,
      "10,224 tokens", size=14, color=ACCENT_DK, bold=True, ls=1.0)
text(s6, wx, col_y + 0.44 + without_h + 0.12, ww, 0.26,
     "Full schema · all 50 tables",
     size=9, color=MUTED, align=PP_ALIGN.CENTER, ls=1.0)

# VS divider
ctext(s6, 5.60, 2.80, 1.10, 0.60, "VS", size=28, color=BORDER_HI, bold=True, ls=1.0)

# With KG column
kx = 6.80; kw = 4.50
text(s6, kx, col_y, kw, 0.34, "With KG",
     size=13, color=KG, bold=True, ls=1.0)
# Bar height proportional: 1655/10224 * 3.20 ≈ 0.52
with_h = 0.52
rect(s6, kx + 0.80, col_y + 0.44, 2.9, with_h, fill=ACCENT, line=None, lw=1.5)
ctext(s6, kx + 0.80, col_y + 0.44 + with_h/2 - 0.14, 2.9, 0.28,
      "1,655 tokens", size=12, color=WHITE, bold=True, ls=1.0)
text(s6, kx, col_y + 0.44 + with_h + 0.12, kw, 0.26,
     "KG-retrieved slice · 8 tables",
     size=9, color=MUTED, align=PP_ALIGN.CENTER, ls=1.0)

# Big stat below
ctext(s6, 0, 5.60, SW, 0.50,
      "83.8% reduction  ·  6.2× factor",
      size=22, color=ACCENT_TX, bold=True, ls=1.0)

footer(s6)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — ReAct Agent Loop
# ─────────────────────────────────────────────────────────────────────────────
s7 = new_slide()
header(s7, "agent", "LLaMA reasons step-by-step via tool calls — max 6 steps")

# Circle loop centred at (5.2, 4.0), radius 1.6
cx_loop = 5.20; cy_loop = 4.00; rad = 1.60
nw = 1.30; nh = 0.60

loop_nodes = [
    # angle (degrees from top), label, fill, line, color
    (90,  "THINK",   T2,      KG,   KG),   # top
    (0,   "ACT",     ACCENT,  None, WHITE), # right
    (270, "OBSERVE", T2,      KG,   KG),   # bottom
    (180, "PLAN",    T2,      KG,   KG),   # left
]

node_centres = []
for ang_deg, lbl, nf, nl, nc in loop_nodes:
    ang = math.radians(ang_deg - 90)  # 90 deg = top
    nx = cx_loop + rad * math.cos(ang) - nw/2
    ny = cy_loop + rad * math.sin(ang) - nh/2
    rect(s7, nx, ny, nw, nh, fill=nf,
         line=nl if nl else nf, lw=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    label_in(s7, nx, ny, nw, nh, lbl, size=11, color=nc)
    node_centres.append((cx_loop + rad * math.cos(ang),
                         cy_loop + rad * math.sin(ang)))

# Connect loop nodes with lines (approximate curved with straight)
for i in range(len(node_centres)):
    x1, y1 = node_centres[i]
    x2, y2 = node_centres[(i+1) % len(node_centres)]
    # offset slightly toward center
    mx = (x1 + x2) / 2; my = (y1 + y2) / 2
    dx = cx_loop - mx; dy = cy_loop - my
    dist = math.sqrt(dx**2 + dy**2)
    if dist > 0:
        ox = dx/dist * 0.3; oy = dy/dist * 0.3
    else:
        ox = 0; oy = 0
    line(s7, x1, y1, x2, y2, color=ACCENT_TX, w=1.5)

# Centre "≤6 steps" label
ctext(s7, cx_loop - 0.55, cy_loop - 0.30, 1.10, 0.60,
      ["≤6", "steps"],
      size=18, color=ACCENT_TX, bold=True, ls=1.0)

# Right panel: Three tool cards
panel_x = 9.30; panel_y = 1.60; panel_w = 3.70
tool_cards = [
    ("get_schema_context", T1, KG,
     "Finds relevant tables + join keys from Neo4j"),
    ("sample_values", SURFACE2, SQL,
     "Resolves categorical filters (status, tier)"),
    ("run_sql", RGBColor(0xE8,0xF5,0xE9), TEAL,
     "Executes read-only SELECT on SQLite"),
]
for i, (tname, tf, tl, tbody) in enumerate(tool_cards):
    ty = panel_y + i * 1.52
    rect(s7, panel_x, ty, panel_w, 1.40, fill=tf, line=tl, lw=1.5)
    text(s7, panel_x+0.12, ty+0.10, panel_w-0.24, 0.28,
         tname, size=9.5, color=tl, bold=True, font=MONO, ls=1.0)
    text(s7, panel_x+0.12, ty+0.44, panel_w-0.24, 0.72,
         tbody, size=9, color=DIM, ls=1.25)

footer(s7)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Tech Stack
# ─────────────────────────────────────────────────────────────────────────────
s8 = new_slide()
header(s8, "technology", "Purpose-built stack — all open-source, all free tier")

tile_cols = cols(3, gap=0.20)
tile_h = 1.50
tile_data = [
    # Row 1
    ("FastAPI",          "REST API · async · Python",           TEAL),
    ("Neo4j AuraDB",     "Graph DB · vector search",            KG),
    ("LLaMA 3.3 70b",    "Groq LPU · 14,400 req/day",          ACCENT_TX),
    # Row 2
    ("fastembed",        "ONNX · all-MiniLM-L6-v2 · 80MB",     SQL),
    ("SQLite",           "50 tables · 23k rows · read-only",    DIM),
    ("Vercel + Render",  "CDN frontend · free API backend",     TEAL),
    # Row 3
    ("python-pptx",      "Deck generation · this slide",        ACCENT_TX),
    ("sentence-tfmrs",   "→ fastembed ONNX (OOM fix)",          SQL),
    ("python-dotenv",    "Secrets · never committed",           TEAL),
]

for i, (val, lbl, col) in enumerate(tile_data):
    row = i // 3; ci = i % 3
    tx, tw = tile_cols[ci]
    ty = 1.55 + row * (tile_h + 0.20)
    tile(s8, tx, ty, tw, tile_h, val, lbl, color=col, vsize=15)

footer(s8)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Security & Guardrails
# ─────────────────────────────────────────────────────────────────────────────
s9 = new_slide()
header(s9, "security", "Three layers prevent data leaks and SQL injection")

card_h = 1.30; gap9 = 0.22; card_start_y = 1.55
card_defs = [
    (TEAL,  "Read-Only SQL",
     "Allowlist of SELECT/WITH only · SQLite opened ?mode=ro · identifier validation strips injection",
     "ENFORCED", TEAL),
    (KG,    "Schema Grounding",
     "KG retriever seeds only real table/column names · model never invents identifiers",
     "ENFORCED", KG),
    (SQL,   "Answer Grounding",
     "Deterministic check: do answer figures appear in the query result rows? Advisory — never gates.",
     "ADVISORY", SQL),
]

for i, (bar_color, title, body, badge_txt, badge_color) in enumerate(card_defs):
    cy = card_start_y + i * (card_h + gap9)
    rect(s9, M, cy, CW, card_h, fill=SURFACE, line=BORDER, lw=1.0)
    rect(s9, M, cy, CW, 0.07, fill=bar_color, line=None)

    # Left colored square
    rect(s9, M+0.12, cy+0.22, 0.55, 0.55, fill=bar_color, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    label_in(s9, M+0.12, cy+0.22, 0.55, 0.55, "✓", size=18, color=WHITE)

    # Title
    text(s9, M+0.82, cy+0.16, CW-3.20, 0.30,
         title, size=13, color=TEXT, bold=True, ls=1.0)
    # Body
    text(s9, M+0.82, cy+0.50, CW-3.20, 0.68,
         body, size=9.5, color=DIM, ls=1.25)

    # Badge
    badge_w = 1.30; badge_h = 0.36
    bx = M + CW - badge_w - 0.10
    by = cy + (card_h - badge_h) / 2
    rect(s9, bx, by, badge_w, badge_h, fill=badge_color, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    label_in(s9, bx, by, badge_w, badge_h, badge_txt, size=8.5, color=WHITE)

footer(s9)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Results Dashboard
# ─────────────────────────────────────────────────────────────────────────────
s10 = new_slide()
header(s10, "results", "Production metrics at a glance")

hero_cols = cols(4, gap=0.20)
hero_row1 = [
    ("50",   "tables in SQLite",     ACCENT_TX),
    ("23k",  "rows of sales data",   SQL),
    ("441",  "KG nodes",             KG),
    ("3–8s", "query response time",  TEAL),
]
hero_row2 = [
    ("83.8%", "context reduction",    ACCENT_TX),
    ("6.2×",  "token efficiency",     KG),
    ("14,400","free queries / day",   SQL),
    ("≤6",    "ReAct steps",          TEAL),
]

for i, (val, lbl, col) in enumerate(hero_row1):
    hx, hw = hero_cols[i]
    hero(s10, hx, 1.55, hw, val, lbl, color=col, vsize=40, lsize=8)

for i, (val, lbl, col) in enumerate(hero_row2):
    hx, hw = hero_cols[i]
    hero(s10, hx, 3.30, hw, val, lbl, color=col, vsize=40, lsize=8)

# Bottom banner
rect(s10, M, 5.20, CW, 0.80, fill=ACCENT, line=None)
ctext(s10, M, 5.30, CW, 0.60,
      "Deployed · Vercel frontend + Render API · Neo4j AuraDB · Groq LPU",
      size=12, color=WHITE, bold=True, ls=1.0)

footer(s10)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — To Do / To Explore
# ─────────────────────────────────────────────────────────────────────────────
s11 = new_slide()
header(s11, "roadmap", "What's done — what's next")

panel_y = 1.55; panel_h = 4.60
lw11 = 5.70; rw11 = 5.36
rx11 = SW - M - rw11  # right panel x

# Left DONE panel
rect(s11, M, panel_y, lw11, panel_h, fill=SURFACE, line=BORDER, lw=1.0)
rect(s11, M, panel_y, lw11, 0.40, fill=TEAL, line=None)
label_in(s11, M, panel_y, lw11, 0.40, "✓  DELIVERED", size=11, color=WHITE)

done_items = [
    "✓  Railway → Render migration (free tier, no timeout)",
    "✓  Groq / LLaMA-3.3-70b integration",
    "✓  sentence-transformers → fastembed (OOM fix)",
    "✓  KG-grounded context slicing (83.8% reduction)",
    "✓  Read-only SQL + identifier guardrails",
    "✓  Vercel CDN + Render API split deployment",
]
for i, item in enumerate(done_items):
    text(s11, M+0.20, panel_y + 0.55 + i*0.66, lw11-0.40, 0.56,
         item, size=10, color=TEAL, ls=1.20)

# Right EXPLORE panel
rect(s11, rx11, panel_y, rw11, panel_h, fill=SURFACE, line=BORDER, lw=1.0)
rect(s11, rx11, panel_y, rw11, 0.40, fill=ACCENT, line=None)
label_in(s11, rx11, panel_y, rw11, 0.40, "?  TO EXPLORE", size=11, color=WHITE)

explore_items = [
    "?  Multi-hop KG joins for complex aggregations",
    "?  Streaming responses (SSE / WebSocket)",
    "?  Authentication layer (OAuth 2.0 / JWT)",
    "?  Auto-refresh KG on schema change",
    "?  Feedback loop — rate answers to fine-tune",
    "?  Multi-database support (Postgres, Snowflake)",
]
for i, item in enumerate(explore_items):
    text(s11, rx11+0.20, panel_y + 0.55 + i*0.66, rw11-0.40, 0.56,
         item, size=10, color=ACCENT_TX, ls=1.20)

footer(s11)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Close (full-bleed gradient, no footer)
# ─────────────────────────────────────────────────────────────────────────────
s12 = new_slide(fill=ACCENT_DK)
grad(s12, 0, 0, SW, SH, ACCENT_DK, ACCENT, angle=135)

# Large "DataPulse"
ctext(s12, 0, 1.80, SW, 1.0, "DataPulse",
      size=62, color=WHITE, bold=True, font=LIGHT, ls=1.0)

# Subtitle
ctext(s12, 0, 2.85, SW, 0.44,
      "Text-to-SQL. Grounded. Fast. Production-ready.",
      size=18, color=WHITE, bold=False, ls=1.0)

# Three stat pills
pills = ["83.8% less context", "3–8s response", "Free tier · Groq + Render"]
pill_w = 2.90; pill_gap = 0.40
total_pill = len(pills)*pill_w + (len(pills)-1)*pill_gap
pill_x_start = (SW - total_pill) / 2
for i, ptxt in enumerate(pills):
    px = pill_x_start + i*(pill_w+pill_gap)
    rect(s12, px, 3.80, pill_w, 0.48, fill=ACCENT_DK, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    label_in(s12, px, 3.80, pill_w, 0.48, ptxt, size=10.5, color=WHITE)

# Contact line
ctext(s12, 0, 5.60, SW, 0.30,
      "Mohan Kumar S  ·  Accenture  ·  mohan.a.kumar.s@accenture.com",
      size=11, color=T2, bold=False, ls=1.0)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
import os
OUT_FINAL = OUT
if os.path.exists(OUT):
    try:
        os.remove(OUT)
    except PermissionError:
        # File may be open in PowerPoint — save with alternate name
        OUT_FINAL = OUT.replace(".pptx", "_v2.pptx")
        if os.path.exists(OUT_FINAL):
            os.remove(OUT_FINAL)
prs.save(OUT_FINAL)
print(f"Saved {OUT_FINAL} — {len(prs.slides)} slides")
